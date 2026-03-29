import os
import boto3
from io import BytesIO
from docx import Document  
import fitz  
import pandas as pd  
from olefile import OleFileIO  
from src.trmeric_api.logging.AppLogger import appLogger
import traceback
from dotenv import load_dotenv
from pptx import Presentation  # Add python-pptx for PPTX support
import tempfile
import docx2txt
from PIL import Image
import base64
from io import BytesIO
import pytesseract




load_dotenv()

AWS_ACCESS_KEY_ID = os.getenv("AWS_ACCESS_KEY_ID")
AWS_SECRET_ACCESS_KEY = os.getenv("AWS_SECRET_ACCESS_KEY")
AWS_S3_BUCKET = os.getenv("AWS_S3_BUCKET")

class S3Service:
    def __init__(self):
        self.s3 = boto3.client(
            's3',
            aws_access_key_id=AWS_ACCESS_KEY_ID,
            aws_secret_access_key=AWS_SECRET_ACCESS_KEY,
            region_name="us-east-1",
        )
        self.bucket_name = AWS_S3_BUCKET
        

    def download_file_as_pd(self, s3_key):
        print(f"--debug S3 key being used: {s3_key}")
        try:
            print("debug 111--- ", s3_key)
            response = self.s3.get_object(Bucket=self.bucket_name, Key=s3_key)
            print("--debug response----------",s3_key)
            binary_content = response['Body'].read()
            print("debug 0--- ", binary_content[:100])
            result = self._extract_xlsx_as_df(binary_content)
            print("debug 1--- Result type:", type(result), "Shape:", result.shape if hasattr(result, 'shape') else "N/A")
            return result
        except Exception as e:
            print(f"Error downloading or processing file {s3_key}: {e}")
            return None
        
    def download_image_binary(self, s3_key: str):
        response = self.s3.get_object(
            Bucket=self.bucket_name,
            Key=s3_key
        )
        return response["Body"].read(), response.get("ContentType", "")

    def generate_presigned_url(self, s3_key, expiry=3600):
        """
        Generates a presigned URL to share an S3 object.
        Args:s3_key (str): S3 object key.expiry (int): URL expiry time in seconds.
        Returns:str: Presigned URL as a string.
        """
        try:
            url = self.s3.generate_presigned_url(
                'get_object',
                Params={
                    'Bucket': self.bucket_name,
                    'Key': s3_key
                },
                ExpiresIn=expiry
            )
            return url
        except Exception as e:
            print(f"Error generating presigned URL for {s3_key}: {e}")
            traceback.print_exc()
            return None


    def download_file_as_text(self, s3_key):
        """
        Downloads a file from S3 and processes its content as text based on file type.
        
        Args:  s3_key (str): The S3 key of the file to download.
            
        Returns:str: Extracted text content from the file, or None if an error occurred.
        """
        try:
            response = self.s3.get_object(Bucket=self.bucket_name, Key=s3_key)
            binary_content = response['Body'].read()
            content_type = response['ContentType']
            
            # print("response -- s3 dowenload -- ", self.bucket_name, s3_key, response, content_type)

            # Determine file format based on extension or content type
            if s3_key.lower().endswith('.pdf') or content_type == 'application/pdf':
                return self.read_pdf_text(binary_content)
            
            # # ---- Images ----
            # elif any(s3_key.lower().endswith(ext) for ext in ('.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.tif')) or content_type.startswith('image/'):
            #     # returns dict with bytes + base64 + dims + format
            #     return self.read_image_file(binary_content)
            
            elif any(s3_key.lower().endswith(ext) for ext in ('.jpg', '.jpeg', '.png', '.webp', '.bmp', '.tiff', '.tif')) or content_type.startswith('image/'):
                # return self.ocr_image_to_text(binary_content)
                return self.describe_image_with_gpt(
                image_binary=binary_content,
                # llm=self.llm,
                detail_level="high",
                purpose="Extract all visible information from the image"
            )

            
            elif s3_key.lower().endswith('.docx') or content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                return self._extract_docx_text(binary_content)
            elif s3_key.lower().endswith('.xlsx') or content_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                return self._extract_xlsx_text(binary_content)
            elif s3_key.lower().endswith('.csv') or content_type == 'text/csv':
                return self._extract_csv_text(binary_content)
            elif s3_key.lower().endswith('.doc') or content_type == 'application/msword':
                return self._extract_doc_text(binary_content)
            elif s3_key.lower().endswith('.txt') or content_type.startswith('text/'):
                return self._decode_text(binary_content)
            elif s3_key.lower().endswith('.pptx') or content_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                return self._extract_pptx_text(binary_content)
            elif s3_key.lower().endswith('.ppt') or content_type == 'application/vnd.ms-powerpoint':
                return self._extract_ppt_text(binary_content)
            else:
                print(f"Unsupported file format for {s3_key}")
                return None

        except Exception as e:
            print(f"Error downloading or processing file {s3_key}: {e}", traceback.format_exc())
            return None
        
    def ocr_image_to_text(self, binary):
        """Extracts text from an image using OCR."""
        try:
            img = Image.open(BytesIO(binary))
            text = pytesseract.image_to_string(img)
            return text.strip()
        except Exception as e:
            print(f"OCR failed: {e}")
            return ""
            
    def describe_image_with_gpt(
        self,
        image_binary: bytes,
        # llm,   # ChatGPTClient injected
        detail_level: str = "high",
        purpose: str = "Explain the image clearly for downstream reasoning"
    ) -> str:
        """
        Uses a GPT vision model to understand and describe an image in detail.

        This is a VISUAL INTERPRETATION layer.
        Output is descriptive context, not factual truth.
        """
        from src.trmeric_ml.llm.Types import ChatCompletion, ModelOptions
        from src.trmeric_ml.llm.models.OpenAIClient import ChatGPTClient
        llm = ChatGPTClient()
        self.log_info = None
        self.modelOptions = ModelOptions(model="gpt-4.1", temperature=0.2, max_tokens=30768)
        try:
            image_b64 = base64.b64encode(image_binary).decode("utf-8")

            system_prompt = f"""
    You are a visual reasoning assistant.

    Your task is to DESCRIBE what is visible in the image.
    Do NOT infer intent or business meaning.
    Do NOT invent components or labels.

    Focus on:
    - Visible components or entities
    - Text and labels (if any)
    - Relationships (arrows, grouping, flow)
    - Overall structure (diagram, architecture, chart, UI)

    Be explicit about uncertainty.

    Detail level: {detail_level}
    Purpose: {purpose}
    """

            description = llm.runVision(
                system_prompt=system_prompt,
                image_base64=image_b64,
                user_instruction="Describe this image in detail.",
                model="gpt-4.1-mini",
                max_tokens=5000,
                temperature=0.1,
            )

            return description.strip() if description else ""

        except Exception as e:
            appLogger.error({
                "function": "describe_image_with_gpt",
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return ""


        
    def _extract_pptx_text(self, binary_content):
        """
        Extracts text from a .pptx file's binary content.
        
        Args:
            binary_content (bytes): Binary content of the PPTX file.
            
        Returns:
            str: Extracted text from all slides.
        """
        try:
            presentation = Presentation(BytesIO(binary_content))
            text = []
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        except Exception as e:
            print(f"Error reading PPTX content: {e}")
            appLogger.error({
                "function": "_extract_pptx_text",
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return ""

    def _extract_ppt_text(self, binary_content):
        """
        Extracts text from a .ppt (PowerPoint 97-2003) file's binary content.
        Note: PPT extraction is limited; consider using external tools like win32com for better results.
        
        Args:
            binary_content (bytes): Binary content of the PPT file.
            
        Returns:
            str: Extracted text or empty string if unsupported.
        """
        try:
            with OleFileIO(BytesIO(binary_content)) as ole:
                if ole.exists('PowerPoint Document'):
                    stream = ole.openstream('PowerPoint Document')
                    content = stream.read()
                    # Simplified text extraction; PPT files are complex and may require external tools
                    return content.decode('utf-8', errors='ignore')
                return ""
        except Exception as e:
            print(f"Error reading PPT content: {e}")
            appLogger.error({
                "function": "_extract_ppt_text",
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return ""


    def upload_file(self, local_file_path, s3_key, bucket_name=None):
        try:
            print("🟡 upload_file CALLED")
            print("Local path:", local_file_path)
            print("Bucket:", bucket_name or self.bucket_name)
            print("Key:", s3_key)

            self.s3.upload_file(
                Filename=local_file_path,
                Bucket=bucket_name or self.bucket_name,
                Key=s3_key
            )

            print("🟢 upload_file SUCCESS")
            return True

        except Exception as e:
            print("🔴 upload_file FAILED")
            print(e)
            return False


    def read_pdf_text(self, binary):
        import fitz
        from PIL import Image
        from io import BytesIO
        import pytesseract

        pdf = fitz.open(stream=binary, filetype="pdf")
        final_text = []

        for i, page in enumerate(pdf):
            print("page", i)

            # ---- STEP 1: block-based detection ----
            blocks = page.get_text("blocks")
            # print("blocks ", blocks)

            text_blocks = [
                b[4].strip()
                for b in blocks
                if b[4] and len(b[4].strip()) > 20
            ]

            if text_blocks:
                print("page normal", i)
                final_text.append("\n".join(text_blocks))
                continue

            # ---- STEP 2: avoid false OCR ----
            images = page.get_images(full=True)
            if not images:
                print("page text-encoded but weird", i)
                raw = page.get_text("rawdict")
                raw_text = " ".join(
                    span.get("text", "").strip()
                    for block in raw.get("blocks", [])
                    for line in block.get("lines", [])
                    for span in line.get("spans", [])
                    if span.get("text", "").strip()
                )
                final_text.append(raw_text)
                continue

            # ---- STEP 3: optimized OCR ----
            print("page OCR", i)
            pix = page.get_pixmap(
                dpi=140,
                colorspace=fitz.csGRAY
            )
            img = Image.open(BytesIO(pix.tobytes("png")))

            ocr_text = pytesseract.image_to_string(
                img,
                config="--psm 6"
            )
            final_text.append(ocr_text)

        # return "\n".join(final_text).strip()
        text = "\n".join(final_text).strip()

        # Data quality signals (VERY important for agents)
        if not text or len(text) < 50:
            return {
                "text": text,
                "status": "low_content",
                "pages": len(pdf),
                "confidence": "low"
            }

        return {
            "text": text,
            "status": "success",
            "pages": len(pdf),
            "confidence": "high"
        }



    def read_image_file(self, binary, max_width=1024):
        """Reads PNG/JPG as bytes, optionally resizes, returns ready-to-send LLM format."""

        # Load image
        img = Image.open(BytesIO(binary))

        # Resize if needed
        if img.width > max_width:
            scale = max_width / img.width
            new_height = int(img.height * scale)
            img = img.resize((max_width, new_height), Image.LANCZOS)

        # Convert back to bytes
        buffer = BytesIO()
        img.save(buffer, format=img.format or "PNG")
        resized_bytes = buffer.getvalue()

        # Base64 for LLM API
        img_b64 = base64.b64encode(resized_bytes).decode("utf-8")

        return {
            # "bytes": resized_bytes,
            "base64": img_b64,
            "width": img.width,
            "height": img.height,
            "format": img.format,
        }




    # def _extract_docx_text(self, binary_content):
    #     """Extracts text from a .docx file's binary content."""
    #     doc = Document(BytesIO(binary_content))
    #     return "\n".join(paragraph.text for paragraph in doc.paragraphs)
    
    def _extract_docx_text(self, binary_content):
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(binary_content)
            tmp.flush()
            text = docx2txt.process(tmp.name)
        return text.strip()

    def _extract_xlsx_text(self, binary_content):
        """Extracts text from an .xlsx file's binary content and returns it in markdown table format."""
        xlsx_data = pd.read_excel(BytesIO(binary_content), sheet_name=None)
        all_text = []
        for sheet_name, sheet in xlsx_data.items():
            all_text.append(f"\n## Sheet: {sheet_name}\n")
            all_text.append(sheet.to_markdown(index=False))
        return "\n".join(all_text)

    def _extract_xlsx_as_df(self, binary_content):
        """Extracts a DataFrame from an .xlsx file's binary content.
        
        Returns a df
        """
        if not binary_content:
            print("--debug binary_content is empty or None")
            return {}
        xlsx_data = pd.read_excel(BytesIO(binary_content), sheet_name=None)
        return xlsx_data

    # best practice says tabled data is best given to llm in a markdown table
    def _extract_csv_text(self, binary_content):
        """Extracts text from a .csv file's binary content and returns it as a valid CSV string."""
        try:
            import pandas as pd
            from io import BytesIO
            csv_data = pd.read_csv(BytesIO(binary_content))
            return csv_data.to_csv(index=False, encoding='utf-8')
        except Exception as e:
            print(f"Error processing CSV content: {e}")
            return None



    def _extract_doc_text(self, binary_content):
        """
        Extracts text from a .doc (Word 97-2003) file's binary content using OleFileIO.
        """
        try:
            with OleFileIO(BytesIO(binary_content)) as ole:
                if ole.exists('WordDocument'):
                    stream = ole.openstream('WordDocument')
                    doc_content = stream.read()
                    # Note: Advanced parsing required for .doc binary formats
                    return doc_content.decode('utf-8', errors='ignore')  # Simplified for demonstration
        except Exception as e:
            print(f"Error reading .doc content: {e}")
        return ""

    def _decode_text(self, binary_content):
        """Decodes binary content as text."""
        try:
            return binary_content.decode('utf-8')
        except UnicodeDecodeError:
            try:
                return binary_content.decode('latin-1')
            except UnicodeDecodeError:
                return binary_content.decode('iso-8859-1', errors='ignore')


    def download_file_as_pd_v2(self, s3_key, filename=None):
        try:
            print(f"--debug downloading S3 key: {s3_key}")
            if filename:
                print(f"--debug using filename for format detection: {filename}")
            response = self.s3.get_object(Bucket=self.bucket_name, Key=s3_key)
            binary_content = response['Body'].read()
            
            if not isinstance(binary_content, bytes):
                print(f"--debug binary_content is not bytes: {type(binary_content)}")
                return None
            
            print(f"--debug binary_content length: {len(binary_content)} bytes")
            
            # Determine file type based on filename if provided, otherwise fall back to s3_key
            file_to_check = filename if filename else s3_key
            
            if file_to_check.lower().endswith('.csv'):
                print(f"--debug detected CSV format from: {file_to_check}")
                df = pd.read_csv(BytesIO(binary_content))
                return {"Sheet1": df}
            elif file_to_check.lower().endswith('.xlsx'):
                print(f"--debug detected XLSX format from: {file_to_check}")
                return self._extract_xlsx_as_df_v2(binary_content)
            else:
                print(f"--debug extension not recognized, trying format detection by content")
                # Try to detect format by attempting to parse as different formats
                
                # Try Excel first (most common for mapping)
                try:
                    print(f"--debug attempting XLSX parsing")
                    result = self._extract_xlsx_as_df_v2(binary_content)
                    if result:
                        print(f"--debug successfully parsed as XLSX")
                        return result
                except Exception as xlsx_error:
                    print(f"--debug XLSX parsing failed: {xlsx_error}")
                
                # Try CSV
                try:
                    print(f"--debug attempting CSV parsing")
                    df = pd.read_csv(BytesIO(binary_content))
                    if not df.empty:
                        print(f"--debug successfully parsed as CSV")
                        return {"Sheet1": df}
                except Exception as csv_error:
                    print(f"--debug CSV parsing failed: {csv_error}")
                
                print(f"--debug all format detection attempts failed")
                return None
        except Exception as e:
            print(f"--debug error downloading or processing file {s3_key}: {str(e)}")
            appLogger.error({
                "event": "download_file_as_pd",
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return None

    def _extract_xlsx_as_df_v2(self, binary_content):
        try:
            excel_file = pd.ExcelFile(BytesIO(binary_content))
            sheets = {sheet: pd.read_excel(excel_file, sheet_name=sheet) for sheet in excel_file.sheet_names}
            print(f"--debug extracted sheets: {list(sheets.keys())}")
            return sheets
        except Exception as e:
            print(f"--debug error extracting Excel file: {str(e)}")
            return None
        
    def download(self,key):
        key = "0041fbda-5160-446a-ae23-e2ab7eb3b89b"

        response = self.s3.get_object(Bucket=self.bucket_name, Key=key)
        binary = response["Body"].read()
        df_dict = pd.read_excel(BytesIO(binary), sheet_name=None)

        for name, df in df_dict.items():
            print(f"Sheet: {name}")
            print(df.head())
        
        return df_dict
