# src/services/super_agent_v1/core/actions.py


from typing import Dict, Optional, List
from src.api.logging.AppLogger import appLogger
import traceback
import re
import base64
from src.database.dao import AgentRunDAO, ProjectsDao, RoadmapDao, IntegrationDao, TenantDao, StatsDao
from src.ml.llm.models.OpenAIClient import ChatGPTClient
from src.ml.llm.Types import ChatCompletion
from src.utils.helper.event_bus import event_bus
from src.s3.s3 import S3Service
from src.utils.helper.file_analyser import FileAnalyzer

import os
from src.utils.helper.decorators import log_function_io_and_time
from src.database.dao import FileDao
from src.utils.types.actions import *
from src.utils.types.getter import *
from ..actions import TrucibleActions, WebDataGetter
from src.database.ai_dao.agent import AIDaoAgentDataGetter
from src.utils.vectorstore.client import TrmericVectorStoreClient
from pathlib import Path
from .style import *


class DataActions:
    def __init__(self, tenant_id: int, user_id: int, agent_name="", session_id="", socketio=None, conversation="", mode = 'research'):
        print("DataActions init ", mode)
        self.tenant_id = tenant_id
        self.user_id = user_id
        self.agent_name = agent_name
        self.session_id = session_id
        self.logInfo = {"tenant_id": tenant_id, "user_id": user_id}
        self.file_analyzer = FileAnalyzer(tenant_id=tenant_id)
        self.s3_service = S3Service()
        self.llm = ChatGPTClient()
        self.socketio = socketio
        self.event_bus = event_bus
        self.conversation = conversation
        self.web_data_getter = WebDataGetter(
            tenant_id=self.tenant_id,
            user_id=self.user_id,
            session_id=session_id,
        )
        self.trucible_actions = TrucibleActions(
            tenant_id=self.tenant_id,
            user_id=self.user_id,
            session_id=session_id,
            agent_name='SuperAgent::Trucible'
        )
        self.vectorstore_client = TrmericVectorStoreClient()
        self.ai_dao_getter = AIDaoAgentDataGetter(
            tenant_id=self.tenant_id,
            user_id=self.user_id,
            session_id=session_id,
            conversation=conversation,
            mode=mode
        )
        self.mode = mode
        self.fn_maps = {
            "read_file_details_with_s3_key": self.read_file_details_with_s3_key,
            "read_image_details_with_s3_key": self.read_image_details_with_s3_key,
        }
        self.fn_maps.update(self.web_data_getter.fn_maps)
        self.fn_maps.update(self.ai_dao_getter.fn_maps)
        self.fn_maps.update(self.trucible_actions.fn_maps)
        self.fn_maps["think_aloud_reasoning"] = self._think_aloud_reasoning_stub
        self.fn_maps["emit_sequential_research_update"] = self._emit_sequential_research_update_stub
        self.fn_maps["emit_research_closure"] = self._emit_research_closure_stub
        self.fn_maps["fetch_trmeric_info_from_vectorstore"] = self.fetch_trmeric_info_from_vectorstore
        self.fn_maps["ask_clarification"] = self.ask_clarification

        self.fn_maps["read_files"] = self.read_file
        self.fn_maps["write_markdown_file"] = self.write_markdown_file
        self.fn_maps["update_section_in_markdown_file"] = self.update_section_in_markdown_file
        self.fn_maps["append_section_in_markdown_file"] = self.append_section_in_markdown_file
        self.fn_maps["merge_and_export_research"] = self.merge_and_export_research

        self.fn_maps["read_html_file"] = self.read_html_file
        self.fn_maps["write_html_file_and_export"] = self.write_html_file_and_export
        self.fn_maps["fetch_files_uploaded_in_session"] = self.fetch_files_uploaded_in_session

        self.fn_maps["fetch_accessible_portfolio_data_using_portfolio_agent"] = self.fetch_fetch_accessible_portfolio_data_using_portfolio_agent
        self.fn_maps["accessible_roadmaps_of_user"] = self.fetch_accessible_roadmaps_of_user
        self.fn_maps["accessible_projects_of_user"] = self.fetch_accessible_projects_of_user
        self.fn_maps["get_available_execution_integrations"] = self.get_available_execution_integrations
        # self.fn_maps["write_and_export_analyst_artifact"] = self.write_and_export_analyst_artifact
        self.fn_maps["fetch_additional_project_execution_intelligence"] = self.fetch_integration_data
        self.fn_maps["freeze_section"] = self.freeze_section
        self.fn_maps["validate_section"] = self.validate_section
        
        
        self.fn_maps["generate_report_doc_after_analysis"]  = self._action_generate_report_doc
        self.fn_maps["generate_html_after_analysis"]        = self._action_generate_html
        self.fn_maps["generate_ppt_after_analysis"]         = self._action_generate_ppt

        self.fn_maps["generate_llm_chart"] = self.generate_llm_chart

        self.fn_maps["total_time_spent_by_customer_on_trmeric"] = self.total_time_spent_by_customer_on_trmeric
        self.fn_maps["get_tenant_knowledge_and_entity_relation_and_volume_stats"] = self.get_tenant_data_volume_stats
     
    
    def generate_llm_chart(self, params: dict = {
        "requirement_focus": "Count projects by portfolio to show distribution",
        "chart_intent": "bar | line | pie | column | area — optional type hint"
    }) -> dict:
        """
        Generates chart JSON using LLM understanding of all data fetched so far.

        Use this action when:
        - User asked for a chart or visualisation
        - A fetch action already ran in this run (data exists in results)
        - BUT no chart artifact (charts_*.json) appeared in the fetch result
        meaning the deterministic chart engine could not construct the chart

        Do NOT call this:
        - Before any fetch action has run (no data to synthesize from)
        - If a chart artifact already exists in execution results
        - More than once per run

        ONE call is sufficient. The LLM reads all accumulated results and
        produces Highcharts-compatible chart JSON uploaded directly to S3.

        params:
            requirement_focus (str):
                What the chart should show.
                Be specific — name the measure and dimension.
                Example: "Count of projects per portfolio, to show distribution"
                Example: "Monthly trend of roadmaps created in 2025"
                Example: "Risk count by severity across active projects"

            chart_intent (str):
                Optional hint about chart type or visual style.
                Example: "bar chart", "line trend", "pie breakdown", "stacked column"
                Leave empty if no preference.
        """
        from src.database.presentation_dao import ChartExportService, LLMChartSynthesizer
        import uuid

        requirement_focus = params.get("requirement_focus", "")
        chart_intent      = params.get("chart_intent", "")

        # Injected by SuperAgent.run() before execute_step is called
        # Contains self.results — all data accumulated in this run
        results = params.get("_injected_results", [])

        synthesizer  = LLMChartSynthesizer(
            tenant_id=self.tenant_id,
            user_id=self.user_id,
            session_id=self.session_id,
        )
        chart_output = synthesizer.synthesize(
            results=results,
            requirement_focus=requirement_focus,
            chart_intent=chart_intent,
        )
        # chart_output = {"charts": [...]}  — Highcharts-compatible

        if not chart_output.get("charts"):
            return {
                "exported": False,
                "reason": "LLM could not derive chart data from accumulated results",
            }

        question_id  = AgentRunDAO.get_latest_run_for_session(
            user_id=self.user_id,
            tenant_id=self.tenant_id,
            session_id=self.session_id,
        )
        base_s3_path = (
            f"exports/{self.tenant_id}/{self.session_id}/{question_id}/"
            if question_id
            else f"exports/{self.tenant_id}/{self.session_id}/"
        )

        s3         = S3Service()
        random_id  = uuid.uuid4().hex
        exporter   = ChartExportService()
        local_path = exporter.export_to_json(chart_output)
        s3_key     = f"{base_s3_path}charts_llm_{random_id}.json"
        s3.upload_file(local_path, s3_key)

        self.event_bus.dispatch(
            "AGENT_ARTIFACT_CREATED",
            {
                "artifact_type": "export",
                "format": "chart",
                # "entity_type": entity_type,
                "requirement_focus": requirement_focus,
                "s3_key": s3_key,
                "source": "llm_synthesized",
            },
            session_id=self.session_id,
        )

        return {
            "exported": True,
            "charts_count": len(chart_output["charts"]),
            "s3_key": s3_key,
            "source": "llm_synthesized",
            "message": "Chart data exported (LLM synthesized)",
        }


    def _action_generate_html(self, params: dict = {
        "html_spec": {
            "title": "",
            "purpose": "executive_report | dashboard | summary | year_review | analysis",
            "audience": "management | executive | engineering | customer",
            "tone": "professional | celebratory | analytical | confident",
            "color_theme": "dark | light | brand",
            "sections_to_cover": ["overview", "metrics", "highlights", "risks"],
            "key_emphasis": "what to highlight most visually",
            "data_personality": "growth | stability | innovation | challenge | mixed"
        }
    }) -> dict:
        """
        Generates a polished single-file HTML document from all analysis results.

        Call this when:
        - All required data has been fetched
        - User requested an HTML report, dashboard, or visual document
        - rough plan artifact_formats contains "html"

        Do NOT call before fetching relevant data.
        ONE call is sufficient — do not repeat.

        params:
            html_spec (dict):
                title (str): document title
                purpose (str): executive_report | dashboard | summary | year_review | analysis
                audience (str): management | executive | engineering | customer
                tone (str): professional | celebratory | analytical | confident
                color_theme (str): dark | light | brand
                sections_to_cover (list): sections to include
                key_emphasis (str): what to highlight most visually
                data_personality (str): growth | stability | innovation | challenge | mixed
        """


    def _action_generate_ppt(self, params: dict = {
        "ppt_spec": {
            "title": "",
            "purpose": "executive_update | project_status | analysis | pitch | review",
            "audience": "management | executive | engineering | board | customer",
            "tone": "professional | energetic | analytical | formal",
            "slide_count_hint": 8,
            "color_theme": "slate_gold | indigo_fire | forest_executive | crimson_modern | ocean_deep | rose_gold",
            "sections_to_cover": ["overview", "highlights", "risks", "next_steps"],
            "key_emphasis": "what to make most prominent in the deck",
            "quality_feedback": []
        }
    }) -> dict:
        """
        Generates a professional PowerPoint (.pptx) from all analysis results collected
        during the current agent run. Uses PptxGenJS via Node.js.

        ────────────────────────────────────────────────────────────
        WHEN TO CALL
        ────────────────────────────────────────────────────────────

        Call this action when:
          - All required data has been fetched in this run
          - User asked for a presentation, deck, slides, or PPT
          - rough_plan artifact_formats contains "ppt"

        Do NOT call:
          - Before fetching relevant data (produce phase must follow fetch)
          - More than once per run (one call is always sufficient)
          - If an artifact already exists in execution results

        ────────────────────────────────────────────────────────────
        PARAMETERS
        ────────────────────────────────────────────────────────────

        params.ppt_spec (dict) — all fields optional, derive from context if absent:

          title (str):
            Presentation title. Keep it short and specific.
            Example: "Q1 2025 Portfolio Review"

          purpose (str):
            The primary intent of the deck.
            Allowed: executive_update | project_status | analysis | pitch | review

          audience (str):
            Who will view this deck. Influences tone, depth, and emphasis.
            Allowed: management | executive | engineering | board | customer

          tone (str):
            Writing and visual tone.
            Allowed: professional | energetic | analytical | formal

          slide_count_hint (int):
            Approximate number of slides to generate. Default: 8.
            Range: 5–12. More than 12 slides degrades quality.

          color_theme (str):
            Visual palette for the deck. Choose based on context:

            slate_gold        → clean, premium, versatile — default choice
                                best for: executive reviews, board decks, portfolio reports

            indigo_fire       → bold, energetic, product-forward
                                best for: product reviews, innovation decks, startup pitches

            forest_executive  → sophisticated, calm, trustworthy
                                best for: finance reviews, strategy docs, investment presentations

            crimson_modern    → bold, urgent, high-stakes
                                best for: risk reviews, operational decks, turnaround presentations

            ocean_deep        → calm, analytical, data-heavy
                                best for: data analysis, quarterly reviews, technical presentations

            rose_gold         → premium, warm, celebratory
                                best for: year reviews, award presentations, customer success decks

            Default: slate_gold (use this when user doesn't specify a color preference)

          sections_to_cover (list[str]):
            Ordered list of topics to include as slides.
            If empty, the generator derives sections from the data and user request.
            Example: ["executive_summary", "project_status", "budget", "risks", "next_steps"]

          key_emphasis (str):
            The single most important thing to make visually prominent.
            Example: "The $420K budget overrun on Lead Gen Automation"
            If empty, derived from the most significant data point in results.

          quality_feedback (list[str]):
            ONLY set this on a RETRY (after quality gate failure).
            Leave empty [] on first call.
            Contains specific issues to fix from the previous attempt.
            Example: ["Stat cards not visible on dark background",
                      "Hero number missing — $420K should dominate slide 4"]

        ────────────────────────────────────────────────────────────
        OUTPUT
        ────────────────────────────────────────────────────────────

        Returns dict:
          exported (bool):      True if .pptx was successfully created and uploaded
          format (str):         "pptx"
          path (str):           "presentation.pptx"
          export_result (dict): S3 upload metadata
          node_stdout (str):    Node.js execution log (useful for debugging)

        ────────────────────────────────────────────────────────────
        DESIGN PRINCIPLES (for the planner to communicate intent)
        ────────────────────────────────────────────────────────────

        The generator automatically:
          - Finds the hero number and makes it the visual surprise slide
          - Maps status signals (on_track/at_risk/compromised) to color
          - Uses sandwich structure (dark title → light content → dark closing)
          - Selects the appropriate skeleton layout per slide type
          - Derives section titles from real data — never placeholders

        The planner should only specify INTENT, not structure.
        Do NOT describe slide layouts or design choices in ppt_spec.
        """


    def _action_generate_report_doc(self, params: dict = {
        "doc_spec": {
            "title": "string",
            "document_type": "prd | strategy | analysis | status_report",
            "audience": "management | executive | engineering",
            "tone": "analytical | formal | concise",
            "scope": "what this document covers",
            "key_emphasis": "what to highlight most"
        }
    }) -> dict:
        """
        Generate a structured document artifact (DOCX) summarizing the analysis results collected
        during the current agent run.

        This action invokes the document writer agent, which converts the analysis outputs
        (data fetched, insights generated, and contextual information) into a professionally
        structured document and exports it as a downloadable artifact.

        The writer agent automatically selects an appropriate document structure based on
        the requested `document_type` and the analysis data available. The planner should
        NOT attempt to design document structure or sections.

        This action should be called **once per request**, after all required data has been
        fetched and analyzed.

        ------------------------------------------------------------------------
        WHEN TO USE
        ------------------------------------------------------------------------

        Call this action when:

        - The user explicitly asks for a document such as:
            - report
            - product requirements document (PRD)
            - strategy document
            - analysis summary
            - status report
            - executive briefing
            - one-pager

        - All required enterprise data has already been fetched and analyzed.

        - The planner determines that a persistent document artifact should be generated
        and exported for the user.

        - The rough plan indicates that an artifact should be produced
        (e.g., `artifact_formats` contains `"report_doc"`).

        ------------------------------------------------------------------------
        WHEN NOT TO USE
        ------------------------------------------------------------------------

        Do NOT call this action when:

        - Required data has not yet been fetched.
        - The user only wants a short conversational answer.
        - A table, spreadsheet, or visualization is more appropriate.
        - Another document has already been generated in the same run.

        Only **one document generation call is needed per request**.

        ------------------------------------------------------------------------
        DOCUMENT TYPES
        ------------------------------------------------------------------------

        The `document_type` field determines the overall structure and writing style.
        The writer agent internally applies an appropriate template for each type.

        Supported document types include:

        - "analysis"
            Analytical report explaining findings, insights, and interpretation
            of the data collected.

        - "summary"
            Concise executive summary highlighting key takeaways.

        - "status_report"
            Progress report covering milestones, risks, team status, and next steps.

        - "prd"
            Product Requirements Document describing objectives, requirements,
            scope, KPIs, and delivery plan.

        - "strategy"
            Strategic roadmap document outlining initiatives, priorities,
            dependencies, and alignment with organizational goals.

        ------------------------------------------------------------------------
        PARAMETERS
        ------------------------------------------------------------------------

        params (dict)

            doc_spec (dict):
                Configuration describing the document to generate.

                title (str)
                    Human-readable document title.

                    Example:
                    "Product Requirements Document: Digital Offerings Scale-Up Accelerator"

                document_type (str)
                    Type of document to generate.

                    Allowed values:
                        "analysis"
                        "summary"
                        "status_report"
                        "prd"
                        "strategy"

                audience (str)
                    Intended audience for the document.

                    Typical values:
                        "management"
                        "executive"
                        "engineering"
                        "customer"

                    This influences writing tone, level of detail,
                    and terminology used.

                tone (str)
                    Writing tone to apply.

                    Common values:
                        "analytical"
                        "formal"
                        "concise"
                        "narrative"

                scope (str)
                    Precise description of what the document covers.

                    Example:
                    "Comprehensive PRD for the Digital Offerings Scale-Up Accelerator roadmap,
                    covering objectives, initiatives, timelines, KPIs, dependencies,
                    constraints, and strategic alignment."

                key_emphasis (str)
                    What aspects should be emphasized in the document.

                    Examples:
                        "Clarity of KPIs and measurable outcomes"
                        "Strategic alignment with digital transformation"
                        "Operational risks and mitigation plans"

        ------------------------------------------------------------------------
        OUTPUT
        ------------------------------------------------------------------------

        Returns:
            dict

            {
                "exported": bool,
                "format": "docx",
                "path": str,
                "export_result": dict
            }

            exported
                True if the document was successfully generated and uploaded.

            format
                Artifact format (currently DOCX).

            path
                Generated document file path.

            export_result
                Metadata returned from the export pipeline (S3 key, artifact info).

        ------------------------------------------------------------------------
        PIPELINE
        ------------------------------------------------------------------------

        Internally, this action performs the following steps:

        1. Compile analysis results and contextual data.
        2. Invoke the document writer agent to produce structured Markdown.
        3. Convert Markdown to DOCX using the document export pipeline.
        4. Upload the artifact to storage (S3).
        5. Emit an artifact creation event for the UI.
        6. Return metadata describing the exported document.

        ------------------------------------------------------------------------
        NOTES
        ------------------------------------------------------------------------

        - The writer agent is responsible for selecting document sections
        and organizing content appropriately.

        - The planner should only specify **document intent**, not structure.

        - The generated document will automatically incorporate:
            - analysis findings
            - enterprise data
            - contextual organization knowledge
            - user request context

        """

      
        
    @log_function_io_and_time
    def get_tenant_data_volume_stats(self, params: Optional[Dict] = {}) -> Dict:
        """
        Returns high-level data volume statistics of the database for the current tenant across
        core execution domains: Projects, Roadmaps, Ideas, and Portfolios.

        Purpose:
        This function helps the agent understand the scale and density of data
        stored in Trmeric so it can:
        - Estimate query cost and data volume before fetching
        - Decide whether to aggregate, sample, or limit scope
        - Avoid heavy operations on large tenants
        - Adapt analysis strategy based on tenant size

        What it returns (per domain):
        - total_tables: Number of tables in the domain
        - total_columns: Total columns across domain tables
        - total_estimated_rows: Estimated total rows (Postgres statistics)
        - total_size_mb: Total storage size
        - avg_rows_per_tenant: Average rows per tenant
        - avg_size_mb_per_tenant: Average storage per tenant

        Domains included:
        - Project
        - Roadmap
        - Idea
        - Portfolio

        Notes:
        - Row counts are estimates (from pg_stat_user_tables), not exact counts.
        - Useful for planning and reasoning, not for precise reporting.
        - This function does NOT fetch business data — only metadata.

        Use this when:
        - The user asks for large-scale analysis
        - The agent is unsure about data size
        - Planning multi-table or full-year analysis
        """

        try:
            stats = StatsDao.GetCombinedStats()
            print("get_tenant_data_volume_stats ", stats)

            # # Optional: add current tenant context
            # tenant_count = stats.get("Project", {}).get("tenant_count", 1)

            # # Estimate current tenant share
            # for domain, data in stats.items():
            #     if not data:
            #         continue

            #     data["estimated_rows_for_current_tenant"] = data.get("avg_rows_per_tenant", 0)
            #     data["estimated_size_mb_for_current_tenant"] = data.get("avg_size_mb_per_tenant", 0)

            return {
                # "tenant_id": self.tenant_id,
                "tenant_data_volume_summary": stats
            }

        except Exception as e:
            appLogger.error({
                "function": "get_tenant_data_volume_stats",
                "error": str(e),
                "traceback": traceback.format_exc(),
                "tenant_id": self.tenant_id,
                "user_id": self.user_id
            })
            return {}

            
    def total_time_spent_by_customer_on_trmeric(self, params= {}):
        """
            Fetch total time spent on trmeric
        """
        return TenantDao.fetch_total_time_spent_in_trmeric(tenant_id=self.tenant_id)
        

    def get_available_execution_integrations(self, params={}):
        """
        Returns the execution integrations (Jira, ADO, GitHub, etc.)
        connected to each project for the current tenant.

        Use this to determine which execution source should be queried
        for a given project before fetching execution data.
        """
        return IntegrationDao.fetchActiveProjectSourcesForTenant(tenant_id=self.tenant_id)
    

    @log_function_io_and_time
    def fetch_integration_data(self, params: Optional[Dict] = {
        "integration_name": "",
        "project_ids": [],
        "user_detailed_query": ""
    }) -> Dict:
        """
        Fetches real-time project execution data from external integrations
        such as Jira, GitHub, or ADO.

        This function retrieves detailed execution metrics for one or more
        projects from the specified integration and filters the results
        based on the user's query.

        Parameters:
        - integration_name: Execution source to query (e.g., "jira", "github", "ado").
        - project_ids: List of project IDs. If empty, all accessible projects
        for the user are considered.
        - user_detailed_query: Natural language query describing which execution
        metrics or signals are required.

        Note:
        If the appropriate integration for a project is unknown,
        call `get_available_execution_integrations` first to resolve
        the correct execution source.
        """
        try:
            params = params or {}
            integration_name = params.get("integration_name")
            project_ids = params.get("project_ids") or []
            if len(project_ids) == 0:
                project_ids = ProjectsDao.FetchAvailableProject(
                    tenant_id=self.tenant_id,
                    user_id=self.user_id,
                )

            is_jira_on_prem = IntegrationDao.is_user_on_prem_jira(self.user_id)
            from src.services.agents.functions.roadmap_analyst import getIntegrationData
            from src.services.integration.helpers.jira_on_prem_getter import fetch_filtered_integration_data
            data = getIntegrationData(
                integration_name=integration_name,
                project_ids=project_ids,
                tenantID=self.tenant_id,
                userID=self.user_id,
            )
            print("is_jira_on_prem ", is_jira_on_prem)
            if (is_jira_on_prem or int(self.tenant_id) in [212, 776]) and (integration_name == "jira" or integration_name == "github"):
                # Flatten out the grouped structure
                ndata = []
                for project_id, integrations in data.items():
                    for item in integrations:
                        if "integration_data" in item:
                            int_data = item["integration_data"].get("data")
                            if int_data:
                                ndata.append(int_data)

                post_proceeded_data = fetch_filtered_integration_data(
                    user_query=params.get("user_detailed_query"),
                    data_array=ndata,
                    integration_name=integration_name
                )
                return post_proceeded_data
            else:
                return data

        except Exception as e:
            appLogger.error({
                "function": "fetch_integration_data",
                "error": str(e),
                "traceback": traceback.format_exc(),
                "tenant_id": self.tenant_id,
                "user_id": self.user_id
            })
            return {}        
        # Consolidated knowledge graph functions (parameter-driven)
        self.fn_maps["fetch_cluster_info"] = self.fetch_cluster_info
        self.fn_maps["fetch_performance_analysis"] = self.fetch_performance_analysis
        
        # Compound analysis functions - high-level combinations for common questions
        self.fn_maps["fetch_performance_landscape"] = self.fetch_performance_landscape
        self.fn_maps["analyze_project_in_context"] = self.analyze_project_in_context
        self.fn_maps["find_success_patterns"] = self.find_success_patterns


    def write_and_export_analyst_artifact(self, params: Dict = {
        "file_name": "analysis_artifact.md",
        "content": "",
        "export_formats": "docx"
    }) -> Dict:
        """
        INTERNAL ACTION — WRITE & EXPORT ANALYST ARTIFACT (COMMUNICATION-ORIENTED)


        Precondition (MANDATORY):
            • params.content MUST be a fully authored markdown document
            • This action MUST NOT be used to “start” writing
            • Writing is assumed complete BEFORE this action is invoked


        Purpose:
            Produce a single, concise, analyst-grade document and immediately
            export it as a user-facing artifact.

            This action represents a COMMUNICATION COMMITMENT, not a research
            commitment. It is intended for decision clarity, explanation,
            or summarization — not deep exploration.

        Core Characteristics:
            • Single document only
            • Explicit export (no auto-export side effects)
            • Deterministic (no rewriting or synthesis beyond what is written)
            • Optimized for clarity, not exhaustiveness

        What this action DOES:
            • Writes one markdown document
            • Immediately exports it (DOCX and/or PDF)
            • Emits an AGENT_ARTIFACT_CREATED event
            • Produces a user-visible deliverable

        What this action DOES NOT do:
            • Does NOT support multiple files
            • Does NOT merge documents
            • Does NOT imply epistemic closure
            • Does NOT replace Deep Research publishing

        When to use:
            • User asks for a report, document, or write-up
            • File-based analysis benefits from a tangible artifact
            • A written deliverable improves decision clarity

        When NOT to use:
            • Multi-dimensional exploration
            • Hypothesis testing
            • Iterative research authoring
            • Situations requiring depth-first investigation
            (escalate to Deep Research instead)

        Parameters:
            params (dict):
                file_name (str, REQUIRED):
                    Name of the output markdown file
                    (e.g., "analysis_summary.md")

                content (str, REQUIRED):
                    Complete markdown content for the document.


        Returns:
            dict:
                {
                    "path": "<markdown_path>",
                    "exported": true,
                }

        Behavioral Contract:
            • This action is an explicit commit
            • The agent must be confident in clarity and correctness
            • This action should be used sparingly and intentionally
        """

        file_name = params.get("file_name")
        content = params.get("content", "")
        export_format = "docx"

        if not file_name:
            return {"error": "file_name is required"}

        if not content.strip():
            return {"error": "content is empty"}

        if not file_name.endswith(".md"):
            file_name = f"{file_name}.md"

        # --------------------------------------------------
        # Write markdown file (single artifact)
        # --------------------------------------------------
        workspace = self.get_workspace()
        full_path = os.path.join(workspace, file_name)
        Path(full_path).parent.mkdir(parents=True, exist_ok=True)

        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content)

        # --------------------------------------------------
        # Explicit export (NO auto-export)
        # --------------------------------------------------
        exports = []
        result = self.export_content({
            "paths": [file_name],
            "export_format": export_format
        })
        exports.append(result)

        return {
            "path": file_name,
            "exported": True,
            "formats": export_format,
            "exports": exports,
        }

    def fetch_fetch_accessible_portfolio_data_using_portfolio_agent(self, params={}):
        """
        This action should be used to fetch basic info of accesible portfolios of user.
        """
        try:
            from src.services.agents import PortfolioApiService
            portfolio_data = PortfolioApiService().get_portfolio_context_of_user(
                user_id=self.user_id,
                tenant_id=self.tenant_id
            )
            return portfolio_data
        except Exception as e:
            return ["Error occured in fetching portfolios info " + str(e)]

    def fetch_accessible_roadmaps_of_user(self, params={}):
        """
        This action should be used to fetch basic info like id and title of accesible roadmaps
        """
        try:
            roadmap_arr = RoadmapDao.fetchEligibleRoadmapList(
                tenant_id=self.tenant_id,
                user_id=self.user_id
            )
            return roadmap_arr
        except Exception as e:
            return ["Error occured in fetching portfolios info " + str(e)]

    def fetch_accessible_projects_of_user(self, params={}):
        """
        This action should be used to fetch basic info like id and title of accesible projects
        """
        try:
            _ids_ = ProjectsDao.FetchAvailableProject(
                tenant_id=self.tenant_id, user_id=self.user_id)
            project_arr = ProjectsDao.fetchProjectIdTitle(
                tenant_id=self.tenant_id,
                project_ids=_ids_
            )
            return project_arr
        except Exception as e:
            return ["Error occured in fetching portfolios info " + str(e)]

    def fetch_files_uploaded_in_session(self, params={}):
        """
        This action should be used to fetch all files uploaded in session
        """
        files = FileDao.FilesUploadedInS3ForSession(self.session_id)
        print("fetch_files_uploaded_in_session", params, files)
        return files

    # =========================================================================
    # COMPOUND ANALYSIS FUNCTIONS
    # High-level functions that combine multiple operations for common questions
    # 
    # ORCHESTRATION PRINCIPLE:
    # These 3 functions handle ~80% of performance/pattern questions in ONE call.
    # Prevents redundant chaining of atomic functions.
    #
    # Key insight: Do NOT call multiple atomic functions when a compound
    # function exists that answers the same question. For example:
    #
    # ❌ BAD:  rank_clusters_by_performance() + fetch_score_performers_analysis()
    # ✅ GOOD: fetch_performance_landscape() (ONE call, same output)
    #
    # ❌ BAD:  fetch_project_cluster_and_score() + fetch_cluster_performance()
    # ✅ GOOD: analyze_project_in_context() (ONE call, MORE context)
    #
    # ❌ BAD:  rank_clusters_by_performance() + loop through fetch_cluster_performance()
    # ✅ GOOD: find_success_patterns() (ONE call, actionable insights)
    # =========================================================================

    def fetch_performance_landscape(self, params: dict = {"top_n": 3, "bottom_n": 3}):
        """
        Get a complete performance landscape: ALL cluster rankings + sample top/bottom performers.

        Executive portfolio overview combining comprehensive cluster analysis with sample projects.
        Returns ALL cluster rankings + limited individual project performers + strategic insights.

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        BEST FOR THESE QUESTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        • "Give me an overview of project performance"
        • "How is our portfolio performing?"
        • "Show me the performance landscape with examples"
        • "What's the big picture on project scores?"
        • "Portfolio performance summary with top/bottom examples"

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETERS (CRITICAL BEHAVIOR)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        - top_n: Number of top individual projects to include (default: 3)
        - bottom_n: Number of bottom individual projects to include (default: 3)
        
        ⚠️ CRITICAL: These parameters ONLY limit individual project performers.
        ⚠️ Cluster rankings ALWAYS include ALL clusters regardless of top_n/bottom_n.
        ⚠️ Example: top_n=3 returns 3 individual projects BUT all 11+ cluster patterns.

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        RETURNS (COMPREHENSIVE)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        {
            "cluster_rankings": [
                {
                    "rank": 1,
                    "pattern_name": "Digital Transformation Pattern",
                    "avg_score": 82.5,
                    "project_count": 5
                },
                ...
            ],
            "top_performers": [
                {
                    "project_title": "Cloud Migration",
                    "core_score": 92,
                    "cluster": {"pattern_name": "Digital Transformation"},
                    "vs_cluster_avg": "+9.5 points above cluster average"
                },
                ...
            ],
            "bottom_performers": [...],
            "insights": {
                "best_cluster": "Digital Transformation (avg 82.5)",
                "worst_cluster": "Legacy Maintenance (avg 58.2)",
                "score_spread": "24 points between best/worst clusters",
                "top_performer_clusters": ["Digital Transformation", "AI/ML"],
                "bottom_performer_clusters": ["Legacy Maintenance"]
            }
        }

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        WHEN TO USE vs OTHER FUNCTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ USE THIS for:
           - Portfolio overview WITH sample project performers
           - Executive summary combining clusters + individual examples
           - "Show performance landscape" (implies clusters + samples)
        
        ❌ DO NOT USE for:
           - "What patterns exist?" → USE fetch_cluster_info(list_all=True)
           - "List all project types" → USE fetch_cluster_info(list_all=True)
           - Just cluster rankings → USE fetch_performance_analysis(analysis_type="rankings")
           - Just top/bottom projects → USE fetch_performance_analysis(analysis_type="performers")
           - Single project → USE analyze_project_in_context
        
        PRIMARY PURPOSE: Comprehensive portfolio overview with both strategic (clusters)
        and tactical (sample performers) perspectives in one call.
        """
        if params is None:
            params = {}

        top_n = params.get("top_n", 3)
        bottom_n = params.get("bottom_n", 3)
        tenant_id = self.tenant_id

        if not tenant_id:
            return {"error": "missing_tenant_id"}

        try:
            appLogger.info({
                "event": "fetch_performance_landscape_start",
                "tenant_id": tenant_id,
                "top_n": top_n,
                "bottom_n": bottom_n
            })
            
            result = _fetch_performance_landscape(
                tenant_id=tenant_id,
                top_n=top_n,
                bottom_n=bottom_n
            )
            
            appLogger.info({
                "event": "fetch_performance_landscape_success",
                "tenant_id": tenant_id
            })
            
            return result

        except Exception as e:
            appLogger.error({
                "event": "fetch_performance_landscape_error",
                "tenant_id": tenant_id,
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return {"error": "performance_landscape_failed", "detail": str(e)}

    def analyze_project_in_context(self, params: dict = {"project_id": ""}):
        """
        Analyze a single project with full context: score, cluster, peer comparison, recommendations.

        Deep-dive into one specific project showing comprehensive performance analysis.
        Includes: core scores + cluster membership + peer ranking + vs cluster average + recommendations.

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        BEST FOR THESE QUESTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        • "Tell me about Project X's performance"
        • "How is [project name] doing compared to similar projects?"
        • "Give me the full picture on this project"
        • "Is Project X above or below average for its type?"
        • "Analyze Project Y in detail"
        • "Where does Project Z rank among its peers?"
        • "Should we be concerned about Project X's performance?"

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETERS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        - project_id: The project ID to analyze (required)
        
        ⚠️ CRITICAL: project_id MUST be the actual project ID (e.g., "4928"), 
        NOT the project title (e.g., "Digital Transformation Initiative").
        
        To find project_id from title, use accessible_projects_of_user first.

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        RETURNS (COMPREHENSIVE)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        {
            "project": {
                "project_title": "Digital Transformation Initiative",
                "core_score": 85,
                "on_time_score": 88,
                "on_scope_score": 82,
                "score_interpretation": "excellent"
            },
            "cluster_context": {
                "pattern_name": "Digital Transformation Pattern",
                "cluster_avg_score": 75.5,
                "project_rank_in_cluster": "1 of 5",
                "vs_cluster_avg": "+9.5 points above average",
                "performance_tier": "top_performer_in_cluster"
            },
            "peer_comparison": {
                "cluster_projects_count": 5,
                "projects_scoring_higher": 0,
                "projects_scoring_lower": 4,
                "peer_projects": [...]
            },
            "recommendations": [
                "This project is a top performer - consider documenting its practices",
                "Team health is strong - could benefit similar projects"
            ]
        }

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        WHEN TO USE vs OTHER FUNCTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ USE THIS for:
           - Single project deep-dive with comprehensive context
           - Project-specific performance questions ("how is Project X doing?")
           - Peer comparison and cluster ranking for ONE project
           - Detailed analysis with actionable recommendations
        
        ❌ DO NOT USE for:
           - Multiple projects → USE fetch_performance_analysis or fetch_performance_landscape
           - Just cluster membership → USE fetch_cluster_info(entity_id=project_id)
           - Just project score → USE fetch_performance_analysis(analysis_type="project")
           - Portfolio overview → USE fetch_performance_landscape
        
        PRIMARY PURPOSE: Comprehensive single-project analysis with peer context and recommendations.
        Most detailed view available for individual project assessment.
        """
        if params is None:
            params = {}

        project_id = params.get("project_id") or ""
        tenant_id = self.tenant_id

        if not project_id:
            return {"error": "missing_project_id"}
        if not tenant_id:
            return {"error": "missing_tenant_id"}

        try:
            appLogger.info({
                "event": "analyze_project_in_context_start",
                "project_id": project_id,
                "tenant_id": tenant_id
            })
            
            result = _analyze_project_in_context(
                project_id=str(project_id),
                tenant_id=tenant_id
            )
            
            appLogger.info({
                "event": "analyze_project_in_context_success",
                "project_id": project_id,
                "tenant_id": tenant_id
            })
            
            return result

        except Exception as e:
            appLogger.error({
                "event": "analyze_project_in_context_error",
                "project_id": project_id,
                "tenant_id": tenant_id,
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return {"error": "project_context_analysis_failed", "detail": str(e)}

    def find_success_patterns(self, params: dict = {"top_n": 3}):
        """
        Identify most successful project patterns AND anti-patterns with strategic insights.

        Strategic planning tool analyzing what works, what doesn't, and why.
        Returns: top success patterns + bottom anti-patterns + score gaps + actionable recommendations.

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        BEST FOR THESE QUESTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        • "What types of projects succeed?"
        • "Which project patterns should we replicate?"
        • "What can we learn from our best performers?"
        • "What project characteristics lead to success?"
        • "Show me success patterns we should follow"
        • "What patterns should we avoid?"
        • "What differentiates our best from our worst projects?"

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETERS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        - top_n: Number of top patterns to analyze (default: 3)

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        RETURNS (COMPREHENSIVE)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        {
            "success_patterns": [
                {
                    "rank": 1,
                    "pattern_name": "Digital Transformation Pattern",
                    "avg_score": 82.5,
                    "success_rate": "80% score above 75",
                    "key_characteristics": {
                        "technologies": ["Cloud", "APIs"],
                        "delivery_themes": ["Agile", "Iterative"]
                    },
                    "top_projects": [
                        {"title": "Cloud Migration", "score": 92}
                    ],
                    "what_makes_it_work": [
                        "Dedicated project management",
                        "Clear milestone-based execution"
                    ]
                },
                ...
            ],
            "anti_patterns": [
                {
                    "pattern_name": "Legacy Maintenance",
                    "avg_score": 58.2,
                    "common_issues": ["Scope creep", "Resource constraints"]
                }
            ],
            "actionable_recommendations": [
                "Prioritize Digital Transformation - 82.5 avg vs 65.3 portfolio",
                "Apply practices: dedicated PM, clear milestones"
            ],
            "portfolio_summary": {
                "total_patterns": 8,
                "avg_portfolio_score": 65.3,
                "patterns_above_average": 3
            }
        }

        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        WHEN TO USE vs OTHER FUNCTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ USE THIS for:
           - Strategic planning questions ("what should we focus on?")
           - Best practice identification and replication
           - Learning from top performers AND avoiding anti-patterns
           - Understanding WHY certain patterns succeed
           - Decision-making about future project types
        
        ❌ DO NOT USE for:
           - "What patterns exist?" → USE fetch_cluster_info(list_all=True)
           - "List all project types" → USE fetch_cluster_info(list_all=True)
           - Just cluster rankings → USE fetch_performance_analysis(analysis_type="rankings")
           - Portfolio overview → USE fetch_performance_landscape
        
        PRIMARY PURPOSE: Strategic insight generation for planning and decision-making.
        Focuses on top_n success patterns and top_n failure patterns with actionable recommendations.
        """
        if params is None:
            params = {}

        top_n = params.get("top_n", 3)
        tenant_id = self.tenant_id

        if not tenant_id:
            return {"error": "missing_tenant_id"}

        try:
            appLogger.info({
                "event": "find_success_patterns_start",
                "tenant_id": tenant_id,
                "top_n": top_n
            })
            
            result = _find_success_patterns(
                tenant_id=tenant_id,
                top_n=top_n
            )
            
            appLogger.info({
                "event": "find_success_patterns_success",
                "tenant_id": tenant_id
            })
            
            return result

        except Exception as e:
            appLogger.error({
                "event": "find_success_patterns_error",
                "tenant_id": tenant_id,
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return {"error": "success_patterns_failed", "detail": str(e)}

    # =========================================================================
    # CONSOLIDATED KNOWLEDGE GRAPH FUNCTIONS
    # Parameter-driven consolidations of 8 atomic functions into 2 smart functions
    # =========================================================================

    def fetch_cluster_info(self, params: dict = {"list_all": False, "entity_id": None, "pattern_id": None, "entity_type": "project"}):
        """
        Universal cluster/pattern lookup - find groups of similar projects or roadmaps.
        
        Clusters are groups of similar entities discovered by machine learning based on
        characteristics like project type, technology stack, delivery approach, etc.
        This function answers "what groups exist?" and "what group is X in?"
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETER BEHAVIOR MATRIX
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        
        list_all=True, entity_type="project"
        → Returns: ALL project clusters with member lists and characteristics
        → Use when: "Show me all project patterns", "What clusters exist?"
        
        list_all=True, entity_type="roadmap"
        → Returns: ALL roadmap clusters with member lists
        → Use when: "What roadmap patterns exist?", "List all roadmap groups"
        
        entity_id="4928", entity_type="project"
        → Returns: Which cluster this specific project belongs to
        → Use when: "What cluster is Project 4928 in?", "Find similar projects to X"
        
        entity_id="xyz", entity_type="roadmap"
        → Returns: Which cluster this specific roadmap belongs to
        → Use when: "What roadmap pattern is XYZ in?"
        
        pattern_id="pattern_0_1353", entity_type="project"
        → Returns: Full details of this specific cluster by ID
        → Use when: "Tell me about pattern_0_1353", "Get cluster details"
        
        pattern_id="Audit Process Transformation Pattern", entity_type="project"
        → Returns: Auto-resolves pattern name to ID, then returns cluster details
        → Use when: "Tell me about the Audit pattern", "Projects in Cloud Migration pattern"
        → Note: Uses LLM to match user-friendly names to actual pattern IDs in graph
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PATTERN NAME RESOLUTION (AUTOMATIC)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        When pattern_id is provided as a user-friendly name instead of an internal ID:
        - System automatically resolves pattern names to graph vertex IDs
        - Works with full names ("Audit Process Transformation Pattern")
        - Works with partial names ("the audit pattern", "cloud migration")
        - Works with descriptions ("pattern about cloud infrastructure")
        - If resolution fails, returns helpful error message with suggestion to use list_all=True
        
        This enables natural follow-up questions like:
        - "Tell me about the Audit Process Transformation Pattern"
        - "Which projects are in the cloud migration cluster?"
        - "Show me the first pattern" (after seeing a list)
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        CRITICAL RULES
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        - Specify EXACTLY ONE of: list_all, entity_id, or pattern_id
        - entity_type determines which graph vertices to query (project vs roadmap)
        - This function ONLY returns cluster membership (no performance scores)
        - For cluster performance, use fetch_performance_analysis(analysis_type="cluster")
        - For finding which clusters perform best, use find_success_patterns()
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        WHEN TO USE vs OTHER FUNCTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ USE fetch_cluster_info: Questions about groupings, patterns, similarity
        ❌ USE fetch_performance_analysis: Need performance scores/rankings
        ❌ USE find_success_patterns: Need strategic insights about what works
        ❌ USE fetch_performance_landscape: Need portfolio overview with scores
        """
        if params is None:
            params = {}
        
        list_all = params.get("list_all", False)
        entity_id = params.get("entity_id")
        pattern_id = params.get("pattern_id")
        entity_type = (params.get("entity_type") or "project").lower()
        tenant_id = self.tenant_id
        
        if not tenant_id:
            return {"error": "missing_tenant_id"}
        
        try:
            return _fetch_cluster_info(
                tenant_id=tenant_id,
                list_all=list_all,
                entity_id=entity_id,
                pattern_id=pattern_id,
                entity_type=entity_type
            )
        except Exception as e:
            appLogger.error({
                "event": "fetch_cluster_info_error",
                "tenant_id": tenant_id,
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return {"error": "cluster_info_failed", "detail": str(e)}

    def fetch_performance_analysis(self, params: dict = {
        "analysis_type": "performers",
        "target_id": None,
        "mode": "both",
        "n": 5,
        "with_insights": False
    }):
        """
        Universal performance analysis with intelligent routing (5-in-1 function).
        
        Swiss Army knife for performance queries - routes to 5 different analyses
        based on analysis_type parameter. Optimized to avoid redundant queries.
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETER BEHAVIOR (analysis_type drives routing)
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        
        analysis_type="performers", mode="top", n=5
        → Top 5 highest scoring projects
        
        analysis_type="performers", mode="bottom", n=3
        → Bottom 3 lowest scoring projects
        
        analysis_type="performers", mode="both", n=5
        → Top 5 AND bottom 5 projects
        
        analysis_type="performers", with_insights=True
        → Performers + cluster rankings context (optimized - ONE query!)
        
        analysis_type="cluster", target_id="pattern_0_1353"
        → Performance summary of all projects in this cluster
        → NOTE: target_id can be pattern_id OR pattern name (auto-resolves!)
        
        analysis_type="rankings"
        → All clusters ranked by average member performance
        
        analysis_type="project", target_id="4928"
        → Specific project's score + cluster + peer comparison
        → ⚠️ CRITICAL: target_id MUST be project_id (NOT project title!)
        
        analysis_type="roadmap", target_id="xyz123"
        → Specific roadmap's score + cluster + peer comparison
        → ⚠️ CRITICAL: target_id MUST be roadmap_id (NOT roadmap title!)
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        CRITICAL: target_id REQUIREMENTS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        - For projects: Use project_id (e.g., "4928"), NOT project title
        - For roadmaps: Use roadmap_id (e.g., "xyz123"), NOT roadmap title
        - For clusters: Use pattern_id (e.g., "pattern_0_1353") OR pattern name
          (pattern names auto-resolve via LLM matching)
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        PARAMETER COMPATIBILITY
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ performers → mode, n, with_insights active
        ✅ cluster → target_id REQUIRED
        ✅ project → target_id REQUIRED
        ✅ rankings → no other params needed
        
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        WHEN TO USE vs OTHER FUNCTIONS
        ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
        ✅ USE this: Single-purpose performance queries, flexible routing
        ❌ USE fetch_performance_landscape: Executive summary (combines everything)
        ❌ USE analyze_project_in_context: Deep-dive on ONE project
        ❌ USE find_success_patterns: Strategic "what works" analysis
        """
        if params is None:
            params = {}
        
        analysis_type = params.get("analysis_type", "performers")
        target_id = params.get("target_id")
        mode = params.get("mode", "both")
        n = params.get("n", 5)
        with_insights = params.get("with_insights", False)
        tenant_id = self.tenant_id
        
        if not tenant_id:
            return {"error": "missing_tenant_id"}
        
        try:
            return _fetch_performance_analysis(
                tenant_id=tenant_id,
                analysis_type=analysis_type,
                target_id=target_id,
                mode=mode,
                n=n,
                with_insights=with_insights
            )
        except Exception as e:
            appLogger.error({
                "event": "fetch_performance_analysis_error",
                "analysis_type": analysis_type,
                "tenant_id": tenant_id,
                "error": str(e),
                "traceback": traceback.format_exc()
            })
            return {"error": "performance_analysis_failed", "detail": str(e)}

    def ask_clarification(self, params={"message": ""}):
        """
            Detail on exact clarification to be asked
        """

    def fetch_trmeric_info_from_vectorstore(
        self,
        params: dict = {
            "query": "",
            "top_k": 5
        }
    ) -> dict:
        """
        INTERNAL ACTION — TRMERIC INFORMATION RETRIEVAL (CUSTOMER SUPPORT)

        Purpose:
            Retrieve Trmeric-specific informational content from
            the internal vector store to help the Customer Success agent:

                • Explain Trmeric’s platform, features, and workflows
                • Answer "what / why / how" questions about Trmeric
                • Provide accurate, company-aligned explanations
                • Avoid generic or hallucinated SaaS descriptions

        This action supports customer conversations such as:
            • "What does Trmeric actually do?"
            • "How does Trmeric handle projects and roadmaps?"
            • "What is Trmeric’s approach to execution intelligence?"
            • "What capabilities does Trmeric provide?"

        This action:
            • Performs semantic retrieval over curated Trmeric documents
            • Returns relevant text snippets and their source context
            • Does NOT perform reasoning, interpretation, or summarization

        This action MUST be used when:
            • The question is specifically about Trmeric
            • The answer depends on internal Trmeric material
            • Accuracy and consistency matter more than speed

        This action MUST NOT be used for:
            • External research or competitor info (use web_search)
            • Metrics, analytics, or dashboards (use AIDAO agents)
            • Guessing product behavior or capabilities

        Parameters:
            params (dict):
                query (str):
                    Natural language description of what the user wants
                    to understand about Trmeric.

                top_k (int):
                    Number of relevant snippets to retrieve (default: 5)

        Returns:
            dict:
                {
                    "matches": [
                        {
                            "chunk": "<text snippet>",
                            "metadata": {
                                "doc_id": "...",
                                "doc_type": "...",
                                "source": "..."
                            }
                        }
                    ],
                    "parent_documents": [
                        {
                            "doc_id": "...",
                            "title": "..."
                        }
                    ]
                }

        Notes:
            • Retrieved content is informational, not authoritative truth
            • Absence of matches does NOT imply absence of capability
            • The agent must explain clearly if information is partial
        """

        self.event_bus.dispatch(
            "STEP_UPDATE",
            {"message": "Fetching Trmeric information"},
            session_id=self.session_id,
        )

        return self.vectorstore_client.query(
            query=params.get("query", ""),
            tenant_id=self.tenant_id,
            collection_name="customer_success",
            top_k=params.get("top_k") or 5,
        )

    def _emit_sequential_research_update_stub(self, params=DEFAULT_PARAMS_FOR_SEQUENTIAL_RESEARCH_UPDATE):
        """
            INTERNAL ACTION — PRESENTATION ONLY (INTERMEDIATE RESEARCH UPDATE)

            Purpose:
            Used to present a focused, intermediate research update to the user
            for a SPECIFIC topic or angle within the broader investigation.

            The goal is to articulate:
                • What is currently understood about the given topic
                • What remains unclear, uncertain, or unexplored for that topic

            When to trigger:
                • Research material relevant to the topic already exists
                • Understanding has improved but is not yet complete
                • Further research is expected after this update

            Required Parameters:
                • topic (string):
                    A clear, human-readable description of the research angle
                    being summarized (e.g., "roadmap feasibility constraints",
                    "data quality limitations", "timeline risks").

            What this action does:
                • Presents accumulated research material related to the topic
                • Clearly separates known vs unknown aspects
                • Helps guide the NEXT research step by making gaps explicit

            What this action does NOT do:
                • Does NOT explore new ideas
                • Does NOT introduce new insights
                • Does NOT perform reasoning or analysis
                • Does NOT conclude the investigation

            Critical Rules:
                • The topic MUST be explicit and non-empty
                • This action MUST be followed by further research steps
            unless confidence in sufficiency is already high

            Note:
                This action does not execute any logic.
                Actual rendering and streaming are handled by the SuperAgent.
                This stub exists solely to communicate intent and constraints
                to the planner.
        """

    def _emit_research_closure_stub(self, params={}):
        """
        INTERNAL ACTION — FINAL PRESENTATION

        Purpose:
        Used to deliver the FINAL research synthesis once
        sufficient confidence has been reached.

        When to trigger:
        • Research material is sufficient
        • No further steps would materially improve understanding
        • Confidence in sufficiency is HIGH

        What this action does:
        • Synthesizes all accumulated research material
        • Produces a clear, decision-ready conclusion
        • Explicitly acknowledges remaining uncertainty (if any)

        Critical rules:
        • This action MUST be the FINAL step in the run
        • This action MUST NOT introduce new facts or reasoning
        • This action MUST NOT trigger any further actions

        Note:
        Actual synthesis and streaming are handled by the SuperAgent.
        This stub exists only to communicate intent to the planner.
        """

    def _think_aloud_reasoning_stub(self, params=DEFAULT_PARAMS_FOR_THINKING) -> Dict:
        """
        INTERNAL ACTION — DIRECTIONAL ONLY

        Purpose:
            Resolve interpretive uncertainty about WHAT should be done next and other details.

        """

    @log_function_io_and_time
    def read_file_details_with_s3_key(self, params: Optional[Dict] = DEFAULT_S3_FILE_PARAMS) -> Dict:
        """
        Reads and analyzes files uploaded in the current session from S3 using provided keys.
        Delegates to FileAnalyzer to handle various file types (TXT, DOC, DOCX, PDF, CSV, Excel).
        """
        self.event_bus.dispatch(
            'STEP_UPDATE',
            {'message': "Reading files"},
            session_id=self.session_id
        )
        return self.file_analyzer.analyze_files(params)

    @log_function_io_and_time
    def read_image_details_with_s3_key(
        self,
        params: Optional[Dict] = DEFAULT_S3_IMAGE_READ_PARAMS
    ) -> Dict:
        """
        Reads image files from S3 and optionally interprets them using GPT Vision.

        Image Handling Rule (CRITICAL):

        If a user uploads an image (including screenshots, TANGO images, or system screenshots)
        and asks to "read", "see", "check", or "look at" it:

        - You MUST execute read_image_details_with_s3_key with mode="read"
        - This performs deterministic OCR only
        - This action is ALWAYS allowed
        - This does NOT require user clarification
        - This does NOT interpret or analyze
        - Vision-based description is optional and must be explicitly requested

        Never ask whether OCR is available — it is available.
        Never assume special pipelines based on image type.


        Modes:
            - "read": OCR only (deterministic)
            - "describe": Vision description only
            - "read_and_describe": OCR + Vision
        """

        params = params or {}
        image_keys = params.get("s3_image_keys_to_read", [])
        mode = params.get("mode") or "read"
        vision_purpose = params.get(
            "vision_purpose") or "Describe the image clearly for downstream reasoning"
        detail_level = params.get("detail_level") or "high"

        # ALLOWED_MODES = {"read", "describe", "read_and_describe"}
        # if mode not in ALLOWED_MODES:
        #     raise ValueError(f"Invalid image read mode: {mode}")

        self.event_bus.dispatch(
            'STEP_UPDATE',
            {'message': f"Reading image files (mode={mode})"},
            session_id=self.session_id
        )

        if not image_keys:
            return {"images": [], "image_count": 0}

        results = []

        for s3_key in image_keys:
            try:
                binary, content_type = self.s3_service.download_image_binary(
                    s3_key)

                image_result = {
                    "file_s3_key": s3_key,
                    "content_type": content_type,
                    "vision_used": True,
                }

                # if mode in ["read", "read_and_describe"]:
                #     ocr_text = self.s3_service.ocr_image_to_text(binary)
                #     image_result["ocr_text"] = ocr_text
                #     image_result["has_text"] = bool(ocr_text.strip())

                # if mode in ["describe", "read_and_describe"]:
                image_b64 = base64.b64encode(binary).decode("utf-8")
                image_result["vision_description"] = self.llm.runVision(
                    system_prompt=self._vision_system_prompt(
                        detail_level=detail_level,
                        purpose=vision_purpose,
                    ),
                    image_base64=image_b64,
                    user_instruction="Describe what you see in this image.",
                )

                results.append(image_result)

            except Exception as e:
                appLogger.error({
                    "function": "read_image_details_with_s3_key",
                    "error": str(e),
                    "traceback": traceback.format_exc(),
                    **self.logInfo
                })
                results.append({
                    "file_s3_key": s3_key,
                    "error": str(e)
                })

        return {
            "images": results,
            "image_count": len(results)
        }

    def _vision_system_prompt(self, detail_level: str, purpose: str) -> str:
        """
        Constructs a strict system prompt for GPT-Vision.

        This prompt ensures the model produces a factual,
        uncertainty-aware visual description without inference
        or business reasoning.
        """
        return f"""
            You are a visual perception assistant.

            Your task is to DESCRIBE what is visible in the image.
            Do NOT infer intent, function, or business meaning.
            Do NOT guess labels, components, or relationships
            unless they are explicitly visible.

            Allowed:
            - Describe shapes, boxes, arrows, lines
            - Transcribe visible text exactly as shown
            - Describe spatial relationships (above, connected to, grouped)
            - Describe flow direction if arrows exist
            - State uncertainty clearly

            Disallowed:
            - Interpretation or analysis
            - Optimization suggestions
            - Assumptions about architecture or domain
            - Conclusions or recommendations

            Detail level: {detail_level}
            Purpose: {purpose}

            Output rules:
            - Plain text only
            - Bullet points allowed
            - No markdown
            - No speculation
        """

    def get_workspace(self) -> str:
        """
        Per-run working directory where the agent writes like a human researcher.
        """
        # base = f"/tmp/agent_runs/{self.session_id}"
        # os.makedirs(base, exist_ok=True)

        root = os.path.abspath(os.path.join(
            os.path.dirname(__file__), "../../../../"))
        base = os.path.join(root, f".cache/agent_runs/{self.session_id}")
        os.makedirs(base, exist_ok=True)
        return base

    def read_file(self, params: Dict = {"paths": [""]}) -> Dict:
        """
        INTERNAL ACTION — READ EXISTING INTERNALLY WRITTEN FILES
        
        INTERNAL ACTION — READ EXISTING INTERNALLY WRITTEN FILES

        Purpose:
            Read the current contents of one or more files that the agent
            has previously written during this session.
            Enables the agent to reason over written material — not just
            memory — and is required before any iterative update or edit.

        When to use:
            • Before updating or rewriting any section of a document
            • When continuing work across multiple research steps
            • When checking what has already been written
            • When assessing document completeness or quality
            • When deciding whether a rewrite or targeted update is appropriate

        Parameters:
            params (dict):
                paths (list of string, REQUIRED):
                    Relative paths to the markdown files inside the agent’s
                    working directory (e.g., ["draft.md", "notes.md"]).

        Guarantees:
            • This action is read-only
            • This action does NOT modify files
        """

        from docx import Document
        from pptx import Presentation

        paths = params.get("paths")
        if not paths:
            return {"error": "path is required"}

        res = []

        self.event_bus.dispatch(
            "LLM_PLAN_UPDATE",
            {
                "info": "step", 
                "content": "Reading Files"
            },
            session_id=self.session_id,
        )

        for path in paths:
            full_path = os.path.join(self.get_workspace(), path)

            if not os.path.exists(full_path):
                res.append({"path": path, "exists": False, "content": ""})
                continue

            content = ""

            try:
                # TEXT FILES
                if path.endswith((".md", ".txt", ".json", ".js", ".html", ".py")):
                    with open(full_path, "r", encoding="utf-8") as f:
                        content = f.read()

                # DOCX FILES
                elif path.endswith(".docx"):
                    doc = Document(full_path)
                    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
                    content = "\n".join(paragraphs)

                # PPTX FILES
                elif path.endswith(".pptx"):
                    prs = Presentation(full_path)
                    slides_text = []

                    for i, slide in enumerate(prs.slides):
                        slide_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_content.append(shape.text)

                        if slide_content:
                            slides_text.append(
                                f"\n--- Slide {i+1} ---\n" + "\n".join(slide_content)
                            )

                    content = "\n".join(slides_text)

                else:
                    content = "[Unsupported file type for text extraction]"

            except Exception as e:
                content = f"[Error reading file: {str(e)}]"

            res.append({
                "path": path,
                "exists": True,
                "content": content
            })

        return res


    def write_markdown_file(self, params: Dict = {"path": "", "section_id": "", "instruction": ""}) -> Dict:
        """
        INTERNAL ACTION — LLM WRITES OR FULLY REWRITE A MARKDOWN DOCUMENT
        Instruction should cover:
            what to analyse from the data
            how to analyse the data - interpretation or exactly paste the determinstic analysis done.
            I say always do interpretation and determinstic both
        
        Document model (CRITICAL):
            • There is NO single “main” document.
            • Each markdown file represents ONE coherent dimension or concern.
            • The agent MAY create MULTIPLE markdown files in parallel.

        Examples:
            • architecture_overview.md
            • payment_routing_design.md
            • provisioning_flow.md
            • sku_governance.md

        Purpose:
            Create a new markdown document or completely replace an existing one with updated content.
            This mirrors how a human researcher rewrites a document
            when structure, framing, or assumptions have changed.

        When to use:
            • Initial document creation
            • Major restructuring or reframing
            • Consolidation after multiple research steps
            • Rewriting multiple sections at once
            • Final cleanup before export

        Parameters:
            params (dict):
                path (string, REQUIRED):
                    Target markdown file path (e.g., "draft.md").
                    
                section_id (string, REQUIRED)
                    Current section_id
                    
                instruction (string, REQUIRED and and very important):
                    what to analyse from the data
                    how to analyse the data - interpretation or exactly paste the determinstic analysis done.
                    I say always do interpretation and determinstic both
        """
        # pass
        return
        # path = params.get("path")
        # content = params.get("content", "")

        # if not path:
        #     return {"error": "path is required"}
        
        
        path = params.get("path")
        instruction = params.get("instruction")

        if not path:
            return {"error": "path is required"}
        if not instruction:
            return {"error": "instruction is required"}

        # ------------------------------------------------------------------
        # LLM PROMPTING
        # ------------------------------------------------------------------
        system_prompt = f"""
            You are a senior researcher and technical writer.

            Your task is to AUTHOR a SINGLE MARKDOWN DOCUMENT.

            DOCUMENT AUTHORING RULES (NON-NEGOTIABLE):
            • This document represents ONE coherent concern or dimension
            • Fully rewrite the document from scratch
            • Do NOT reference other documents unless explicitly instructed
            • Do NOT merge unrelated concerns
            • Structure the document clearly with headings and subheadings
            • Be precise, analytical, and decision-oriented
            • Markdown only — no code blocks unless required by content
            • No explanations about what you are doing

            CONTEXT:
            {self.context_string}

            AGENT ROLE:
            {self.policy.agent_role}
        """

        user_prompt = f"""
        You must write the markdown file according to the following instruction.

        INSTRUCTION (AUTHORITATIVE):
        {instruction}

        OUTPUT RULES:
        • Output ONLY valid Markdown
        • No commentary
        • No JSON
        • No placeholders like "TBD"
        • If a section is required, WRITE it fully
        """

        chat = ChatCompletion(
            system=system_prompt,
            prev=[],
            user=user_prompt,
        )

        # ------------------------------------------------------------------
        # STREAM + COLLECT
        # ------------------------------------------------------------------
        content = ""
        for chunk in self.llm.runWithStreaming(
            chat,
            self.model_options,
            "super::write_markdown",
            logInDb=self.log_info
        ):
            content += chunk

        # ------------------------------------------------------------------
        # WRITE FILE
        # ------------------------------------------------------------------
        full_path = os.path.join(self.get_workspace(), path)
        Path(full_path).parent.mkdir(parents=True, exist_ok=True)

        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content.strip())

        # ------------------------------------------------------------------
        # EMIT ARTIFACT EVENT (IMPORTANT FOR UI / TRACEABILITY)
        # ------------------------------------------------------------------
        event_bus.dispatch(
            "AGENT_ARTIFACT_CREATED",
            {
                "artifact_type": "export",
                "path": path,
                "description": "Markdown document authored by LLM"
            },
            session_id=self.session_id
        )

        return {
            "path": path,
            "written": True
        }

        full_path = os.path.join(self.get_workspace(), path)
        Path(full_path).parent.mkdir(parents=True, exist_ok=True)

        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content)

        return {"path": path, "written": True}

    def update_section_in_markdown_file(self, params: Dict = {
            "path": "",
            "section_header": "",
            "new_content": ""
        }
    ) -> Dict:
        """
        INTERNAL ACTION — TARGETED SECTION UPDATE

        Purpose:
            Replace the content of a specific markdown section
            while preserving the rest of the document.

            This enables focused, incremental improvement of
            individual dimensions without destabilizing the full draft.

        When to use:
            • Refining a single analytical or conceptual dimension
            • Incorporating new evidence into a known section
            • Correcting errors or assumptions in one section
            • Improving clarity or rigor locally

        When NOT to use:
            • When multiple sections are affected
            • When document framing or structure has changed
            • When the section does not yet exist
            • When a full rewrite would be cleaner

        Parameters:
            params (dict):
                path (string, REQUIRED):
                    Markdown file path.

                section_header (string, REQUIRED):
                    Exact markdown header identifying the section
                    (e.g., "## Risk Analysis", "### Data Limitations").

                new_content (string, REQUIRED):
                    Full replacement content for that section.
                    This content should be complete and self-contained.

        Returns:
            dict:
                {
                    "path": string,
                    "section_updated": string
                }

        Critical rules:
            • The section header MUST already exist
            • Content is replaced, not appended
            • If more than one section needs changes,
            the agent SHOULD prefer write_markdown_file
        """

        path = params.get("path")
        header = params.get("section_header")
        new_content = params.get("new_content", "")

        if not path or not header:
            return {"error": "path and section_header are required"}

        full_path = os.path.join(self.get_workspace(), path)

        if not os.path.exists(full_path):
            return {"error": "file does not exist"}

        with open(full_path, "r", encoding="utf-8") as f:
            text = f.read()

        pattern = rf"({re.escape(header)}\n)(.*?)(\n## |\Z)"
        match = re.search(pattern, text, flags=re.DOTALL)

        if not match:
            return {"error": "section not found"}

        updated = (
            text[:match.start(2)]
            + new_content.strip() + "\n"
            + text[match.end(2):]
        )

        with open(full_path, "w", encoding="utf-8") as f:
            f.write(updated)

        # self._auto_export_content(path)
        # self.event_bus.dispatch(
        #     "STEP_UPDATE",
        #     {"message": f"Updating section {header}"},
        #     session_id=self.session_id,
        # )

        return {"path": path, "section_updated": header}

    def append_section_in_markdown_file(self, params: Dict = {"path": "notes.md", "content": ""}) -> Dict:
        """
        INTERNAL ACTION — RESEARCH SCRATCHPAD (NON-FINAL)

        Purpose:
            Append informal notes, partial insights, open questions,
            or weak signals discovered during research.

            This acts as the agent’s working notebook and does NOT
            represent polished or final content.

        When to use:
            • Capturing incomplete or tentative insights
            • Recording unanswered questions
            • Parking ideas not yet ready for the main document
            • Tracking contradictions or uncertainties

        When NOT to use:
            • Writing polished sections
            • Updating the main document
            • Final synthesis or conclusions

        Parameters:
            params (dict):
                path (string, OPTIONAL):
                    Notes file path (default: "notes.md").

                content (string, REQUIRED):
                    A single note or observation.
                    Should be concise and informal.

        Returns:
            dict:
                {
                    "path": string,
                    "appended": true
                }

        Behavioral notes:
            • Notes are additive and non-destructive
            • Notes may later be promoted into the main document
            • Notes do NOT imply correctness or completeness
        """

        path = params.get("path", "notes.md")
        content = params.get("content", "").strip()

        if not content:
            return {"error": "content is empty"}

        full_path = os.path.join(self.get_workspace(), path)
        Path(full_path).parent.mkdir(parents=True, exist_ok=True)

        with open(full_path, "a", encoding="utf-8") as f:
            f.write(f"\n- {content}\n")

        # self._auto_export_content(path)
        # self.event_bus.dispatch(
        #     "STEP_UPDATE",
        #     {"message": f"Appending notes to {path}"},
        #     session_id=self.session_id,
        # )

        return {"path": path, "appended": True}

    def merge_and_export_research(self, params: Dict = {"ordered_files": [], "output_path": ""}) -> Dict:
        """
        INTERNAL ACTION — MERGE & EXPORT RESEARCH ARTIFACT (AUTHORITATIVE)

        Purpose:
            Produce the FINAL, authoritative research deliverable by
            deterministically merging multiple markdown research files
            into a single consolidated document and exporting it.

            This action represents a **commit point** in the research lifecycle.
            It should be used ONLY after research is deemed sufficient.

        Core Principles:
            • Deterministic (no rewriting, no synthesis, no inference)
            • Lossless (all content is preserved as written)
            • Order-sensitive (caller explicitly defines merge order)
            • Produces ONE canonical export

        What this action DOES:
            • Reads multiple existing markdown files from the agent workspace
            • Concatenates them sequentially in the provided order
            • Inserts clear section boundaries between files
            • Writes a single merged markdown file
            • Exports the merged file (DOCX / PDF via existing export pipeline)
            • Emits a single AGENT_ARTIFACT_CREATED event

        What this action DOES NOT do:
            • Does NOT rewrite, summarize, or rephrase content
            • Does NOT resolve redundancy or conflicts
            • Does NOT introduce new structure or conclusions
            • Does NOT perform synthesis (that is a separate concern)

        When to use:
            • After confidence_in_sufficiency is HIGH
            • When the user requests a “final document”, “report”, or “deliverable”
            • When research exploration is complete and conclusions are stable

        When NOT to use:
            • During exploration or hypothesis testing
            • While research files are still evolving
            • For executive summarization or narrative polishing

        Parameters:
            params (dict):
                ordered_files (List[str], REQUIRED):
                    A sequential list of markdown file paths (relative to the
                    agent workspace) defining the EXACT merge order.

                    Example:
                        [
                            "architecture_overview.md",
                            "risk_analysis.md",
                            "operational_constraints.md",
                            "open_questions.md"
                        ]

                output_path (str, OPTIONAL):
                    Path for the merged markdown file.
                    Default: "merged_research.md"

                export_format (str, OPTIONAL):
                    Export format for the merged artifact.
                    Default: "docx"
                    Allowed: "docx", "pdf", "md"

        Returns:
            dict:
                {
                    "merged_file": "<output_path>",
                    "files_merged": [ ... ],
                    "exported": true
                }

        Behavioral Contract:
            • Caller is responsible for deciding merge order
            • This action assumes all listed files already exist
            • Failure to read any file should fail the entire action
        """
        try:
            # --------------------------------------------------
            # Validate input
            # --------------------------------------------------
            ordered_files: List[str] = params.get("ordered_files")
            output_path: str = params.get("output_path", "merged_research.md")

            if not ordered_files or not isinstance(ordered_files, list):
                return {"error": "ordered_files (list[str]) is required"}

            if not output_path.endswith(".md"):
                return {"error": "output_path must be a markdown (.md) file"}

            workspace = self.get_workspace()
            merged_sections = []

            # --------------------------------------------------
            # Read & merge files deterministically
            # --------------------------------------------------
            for idx, rel_path in enumerate(ordered_files, start=1):
                full_path = os.path.join(workspace, rel_path)

                if not os.path.isfile(full_path):
                    return {
                        "error": "file_not_found",
                        "missing_file": rel_path
                    }

                with open(full_path, "r", encoding="utf-8") as f:
                    content = f.read().strip()

                # Clear, explicit boundary between files
                section_header = (
                    f"\n\n"
                    f"---\n"
                    # f"## Section {idx}: {os.path.basename(rel_path)}\n"
                    f"---\n\n"
                )
                merged_sections.append(section_header + content)

            merged_content = "\n\n".join(merged_sections).strip() + "\n"

            # --------------------------------------------------
            # Write merged markdown file
            # --------------------------------------------------
            merged_full_path = os.path.join(workspace, output_path)
            Path(merged_full_path).parent.mkdir(parents=True, exist_ok=True)

            with open(merged_full_path, "w", encoding="utf-8") as f:
                f.write(merged_content)

            # --------------------------------------------------
            # Export merged file (DOCX + PDF)
            # --------------------------------------------------
            exports = []

            # DOCX (this will also generate PDF via your existing logic)
            docx_result = self._auto_export_content(output_path, "docx")
            exports.append(docx_result)

            # --------------------------------------------------
            # Final response
            # --------------------------------------------------
            return {
                "merged_file": output_path,
                "files_merged": ordered_files,
                "exported": True,
                "exports": exports
            }
        except Exception as e:
            appLogger.error({
                "function": "merge_and_export_research",
                "error": str(e),
                **self.logInfo
            })
            return {
                "error": str(e),
            }

    def _auto_export_content(self, path, fmt='docx'):
        """
        INTERNAL — auto-export markdown after any mutation.
        This makes every document change immediately user-visible.
        """
        try:
            self.export_content({
                "paths": [path],
                "export_format": fmt
            })

        except Exception as e:
            appLogger.error({
                "function": "_auto_export_content",
                "error": str(e),
                **self.logInfo
            })
            return {
                "error": str(e)
            }

    def export_content(self, params: Dict) -> Dict:
        import os
        import uuid
        import tempfile
        import pypandoc
        from typing import List

        print("export_content ", params)

        # --------------------------------------------------
        # Validate input
        # --------------------------------------------------
        paths: List[str] = params.get("paths")
        export_format: str = params.get("export_format", "md")

        if not paths or not isinstance(paths, list):
            return {"error": "paths (list of markdown files) is required"}

        if export_format not in {"md", "docx", "html", "pptx", "docx_file"}:
            return {"error": f"unsupported export_format: {export_format}"}

        # Ensure pandoc exists
        try:
            pypandoc.get_pandoc_version()
        except OSError:
            pypandoc.download_pandoc()

        workspace = self.get_workspace()
        exports = []

        # --------------------------------------------------
        # Process each markdown file
        # --------------------------------------------------
        for rel_path in paths:
            docx_path = None
            pdf_path = None
            full_path = os.path.join(workspace, rel_path)

            if not os.path.isfile(full_path):
                return {"error": f"file not found: {rel_path}"}

            # with open(full_path, "r", encoding="utf-8") as f:
            #     md_content = f.read()
                
            # only read as text for text-based formats
            if export_format in {"md", "docx", "html"}:
                with open(full_path, "r", encoding="utf-8") as f:
                    md_content = f.read()
            else:
                md_content = None

            normalized_format = "docx" if export_format == "docx_file" else export_format
            file_id = f"{uuid.uuid4()}.{normalized_format}"

            question_id = AgentRunDAO.get_latest_run_for_session(
                user_id=self.user_id,
                tenant_id=self.tenant_id,
                session_id=self.session_id
            )

            base_s3_path = (
                f"exports/{self.tenant_id}/{self.session_id}/{question_id}/"
                if question_id else ""
            )

            random_id = uuid.uuid4().hex
            # export_type = "html" if export_format == "html" else "doc"
            # s3_key = f"{base_s3_path}_{export_type}_{random_id}.{normalized_format}"
            
            export_type = "html" if export_format == "html" else "ppt" if export_format == "pptx" else "doc"
            s3_key = f"{base_s3_path}_{export_type}_{random_id}.{normalized_format}"


            # --------------------------------------------------
            # Temp file
            # --------------------------------------------------
            with tempfile.NamedTemporaryFile(
                delete=False,
                suffix=f".{normalized_format}"
            ) as tmp:
                tmp_path = tmp.name

            try:
                # --------------------------------------
                # RAW MARKDOWN EXPORT
                # --------------------------------------
                if export_format == "md":
                    with open(tmp_path, "w", encoding="utf-8") as f:
                        f.write(md_content)

                    self.s3_service.upload_file(tmp_path, s3_key)

                # --------------------------------------
                # HTML EXPORT
                # --------------------------------------
                elif export_format == "html":
                    with open(tmp_path, "w", encoding="utf-8") as f:
                        f.write(md_content)

                    self.s3_service.upload_file(tmp_path, s3_key)
                    
                    # EMIT EVENT
                    self.event_bus.dispatch(
                        "AGENT_ARTIFACT_CREATED",
                        {
                            "artifact_type": "export",
                            "format": "html",
                            "file_name": file_id,
                            "s3_key": s3_key,
                            "source_file": rel_path,
                        },
                        session_id=self.session_id,
                    )

                # --------------------------------------
                # DOCX + PDF EXPORT
                # --------------------------------------
                elif export_format == "docx":
                    # --------------------------------------------------
                    # Convert MD → HTML (authoritative format)
                    # --------------------------------------------------
                    html_content = pypandoc.convert_text(
                        md_content,
                        to="html",
                        format="markdown+pipe_tables+grid_tables+multiline_tables",
                        extra_args=["--standalone"]
                    )
                    html_content = inject_css(html_content, TABLE_CSS)

                    # ---- DOCX (HTML → DOCX) ----
                    docx_path = tmp_path + ".docx"
                    pypandoc.convert_text(
                        html_content,
                        to="docx",
                        format="html",
                        outputfile=docx_path,
                        extra_args=["--standalone"]
                    )

                    self.s3_service.upload_file(docx_path, s3_key)

                    self.event_bus.dispatch(
                        "AGENT_ARTIFACT_CREATED",
                        {
                            "artifact_type": "export",
                            "format": "docx",
                            "file_name": file_id,
                            "s3_key": s3_key,
                            "source_file": rel_path,
                        },
                        session_id=self.session_id,
                    )

                    # ---- PDF (HTML → PDF) ----
                    pdf_path = tmp_path + ".pdf"
                    pdf_s3_key = s3_key.replace(".docx", ".pdf").replace("_doc_", "_pdf_")

                    pypandoc.convert_text(
                        html_content,
                        to="pdf",
                        format="html",
                        outputfile=pdf_path,
                        extra_args=[
                            "--standalone",
                            "--pdf-engine=wkhtmltopdf",
                            "-V", "encoding=utf-8",
                            "-V", "margin-top=10mm",
                            "-V", "margin-bottom=10mm",
                            "-V", "margin-left=12mm",
                            "-V", "margin-right=12mm",
                            "-V", "page-size=A4",
                            "-V", "dpi=300",
                            "--metadata", "title=Document"
                        ]
                    )

                    self.s3_service.upload_file(pdf_path, pdf_s3_key)

                    self.event_bus.dispatch(
                        "AGENT_ARTIFACT_CREATED",
                        {
                            "artifact_type": "export",
                            "format": "pdf",
                            "file_name": file_id,
                            "s3_key": pdf_s3_key,
                            "source_file": rel_path,
                        },
                        session_id=self.session_id,
                    )

                    if os.path.exists(pdf_path):
                        os.unlink(pdf_path)

                # --------------------------------------
                # PPTX EXPORT (already generated file)
                # --------------------------------------
                elif export_format == "pptx":
                    if not os.path.isfile(full_path):
                        return {"error": f"pptx file not found: {rel_path}"}

                    # ── Upload PPTX ──
                    self.s3_service.upload_file(full_path, s3_key)

                    self.event_bus.dispatch(
                        "AGENT_ARTIFACT_CREATED",
                        {
                            "artifact_type": "export",
                            "format": "pptx",
                            "file_name": file_id,
                            "s3_key": s3_key,
                            "source_file": rel_path,
                        },
                        session_id=self.session_id,
                    )
                # --------------------------------------
                # PRE-BUILT DOCX (generated by docx.js — upload directly, no pandoc)
                # --------------------------------------
                elif export_format == "docx_file":
                    if not os.path.isfile(full_path):
                        return {"error": f"docx file not found: {rel_path}"}
                    
                    self.s3_service.upload_file(full_path, s3_key)
                    self.event_bus.dispatch(
                        "AGENT_ARTIFACT_CREATED",
                        {
                            "artifact_type": "export",
                            "format": "docx",
                            "file_name": file_id,
                            "s3_key": s3_key,
                            "source_file": rel_path,
                        },
                        session_id=self.session_id,
                    )

                    # # ── Convert PPTX → PDF via LibreOffice and upload ──
                    # try:
                    #     import subprocess
                    #     import tempfile

                    #     pdf_s3_key = s3_key.replace("_ppt_", "_pdf_").replace(".pptx", ".pdf")
                    #     pptx_stem  = os.path.splitext(os.path.basename(full_path))[0]

                    #     with tempfile.TemporaryDirectory() as tmp_pdf_dir:
                    #         pdf_path = os.path.join(tmp_pdf_dir, f"{pptx_stem}.pdf")

                    #         result = subprocess.run(
                    #             [
                    #                 "soffice",
                    #                 "--headless",
                    #                 "--norestore",
                    #                 "--nofirststartwizard",
                    #                 "--convert-to", "pdf",
                    #                 "--outdir", tmp_pdf_dir,
                    #                 full_path,
                    #             ],
                    #             capture_output=True,
                    #             text=True,
                    #             timeout=60,
                    #             env={**os.environ, "HOME": "/root"},
                    #         )

                    #         if result.returncode != 0:
                    #             print("LibreOffice PDF conversion failed:", result.stderr)
                    #         elif os.path.isfile(pdf_path):
                    #             pdf_file_id = f"{uuid.uuid4()}.pdf"
                    #             self.s3_service.upload_file(pdf_path, pdf_s3_key)

                    #             self.event_bus.dispatch(
                    #                 "AGENT_ARTIFACT_CREATED",
                    #                 {
                    #                     "artifact_type": "export",
                    #                     "format": "pdf",
                    #                     "file_name": pdf_file_id,
                    #                     "s3_key": pdf_s3_key,
                    #                     "source_file": rel_path,
                    #                 },
                    #                 session_id=self.session_id,
                    #             )

                    #             exports.append({
                    #                 "source_file": rel_path,
                    #                 "file_name": pdf_file_id,
                    #                 "s3_key": pdf_s3_key,
                    #             })
                    #         else:
                    #             print("LibreOffice ran but PDF not found at:", pdf_path)
                    #             print("soffice stdout:", result.stdout)
                    #             print("soffice stderr:", result.stderr)

                    # except subprocess.TimeoutExpired:
                    #     print("LibreOffice PDF conversion timed out")
                    # except Exception as e:
                    #     print("PDF conversion error:", e)
                    # # TemporaryDirectory cleans up automatically — no finally needed
                    
            except Exception as e:
                print("export_content error:", e)
                raise

            # finally:
            #     if os.path.exists(tmp_path):
            #         os.unlink(tmp_path)
            #     if export_format == "docx" and os.path.exists(docx_path):
            #         os.unlink(docx_path)
            finally:
                if os.path.exists(tmp_path):
                    os.unlink(tmp_path)
                if docx_path and os.path.exists(docx_path):
                    os.unlink(docx_path)
                if pdf_path and os.path.exists(pdf_path):
                    os.unlink(pdf_path)

            exports.append({
                "source_file": rel_path,
                "file_name": file_id,
                "s3_key": s3_key,
            })

        return {
            "exported": True,
            "format": export_format,
            "exports": exports,
        }


    def read_html_file(self, params: Dict = {"path": ""}) -> Dict:
        """
        INTERNAL ACTION — READ EXISTING HTML PROTOTYPE

        Purpose:
            Read an HTML prototype previously written by the agent.
            Used for iterative ideation and refinement.

        When to use:
            • Before modifying an existing prototype
            • When reasoning about what is already visualized
            • When iterating collaboratively with the user

        Parameters:
            params (dict):
                path (string, REQUIRED):
                    Relative HTML file path inside the workspace
                    (e.g., "prototype.html", "landing_v1.html")

        Returns:
            dict:
                {
                    "exists": boolean,
                    "path": string,
                    "content": string
                }

        Guarantees:
            • Read-only
            • No parsing or mutation
            • If file does not exist, exists=false
        """

        path = params.get("path")
        if not path:
            return {"error": "path is required"}

        full_path = os.path.join(self.get_workspace(), path)

        if not os.path.exists(full_path):
            return {
                "exists": False,
                "path": path,
                "content": ""
            }

        with open(full_path, "r", encoding="utf-8") as f:
            content = f.read()

        return {
            "exists": True,
            "path": path,
            "content": content
        }

    def write_html_file_and_export(
        self,
        params: Dict = {"path": "", "html_content": ""}
    ):
        """
        INTERNAL ACTION — WRITE OR FULLY REWRITE HTML PROTOTYPE

        HTML authoring model (STRICT):

        • One HTML file = one prototype
        • Entire file is rewritten on every update
        • No partial DOM edits
        • No assumptions of production readiness

        Purpose:
            Materialize an idea as a lightweight HTML prototype
            for visualization and feedback.

        When to use:
            • Creating a new prototype
            • Updating an existing prototype based on feedback
            • Iterating on layout, flow, or interaction ideas

        When NOT to use:
            • For markdown documents
            • For production-ready UI
            • For backend or business logic

        Parameters:
            params (dict):
                path (string, REQUIRED):
                    HTML file path (e.g., "prototype.html")

                html_content (string, REQUIRED):
                    Full HTML content.
                    Must include complete HTML structure.

        Critical rules:
            • File is fully overwritten
            • This action creates a user-visible artifact
            • HTML is treated as the final artifact (no conversion)
        """

        path = params.get("path")
        content = params.get("html_content", "")
        
        print("write_html_file_and_export ", path)

        if not path:
            return {"error": "path is required"}

        if not path.lower().endswith(".html"):
            return {"error": "HTML prototypes must use .html extension"}

        if not content.strip():
            return {"error": "content is empty"}

        print("write_html_file_and_export ", path)
        
        # Basic safety check (non-blocking, ideation-friendly)
        if "<html" not in content.lower():
            appLogger.error({
                "function": "write_html_file_and_export",
                "warning": "HTML content does not include <html> tag",
                **self.logInfo
            })

        full_path = os.path.join(self.get_workspace(), path)
        Path(full_path).parent.mkdir(parents=True, exist_ok=True)
        with open(full_path, "w", encoding="utf-8") as f:
            f.write(content)

        self._auto_export_content(path, "html")
        return {
            "path": path,
            "written": True
        }


    def freeze_section(self, params: Dict = {"section_id": ""}):
        """
        Freeze a research section by mutating the in-memory execution plan state.

        PURPOSE
        -------
        Marks a section as frozen once it has been authored and validated.
        A frozen section is considered epistemically finalized and must not
        be rewritten or re-analyzed within the same run.

        This function operates as a workflow state transition mechanism,
        moving a section from:
            written + validated → frozen

        PRECONDITIONS
        --------------
        • An execution plan must exist
        • The latest plan must contain `current_section_state`
        • The section_id must match the active section
        • Section must be both:
            - written = True
            - validated = True

        SIDE EFFECTS
        -------------
        • Updates:
            current_section_state.frozen = True
            current_section_state.frozen_at = <UTC timestamp>

        • Mutates the latest execution plan in memory
        • Does NOT persist external state
        • Future planning steps inherit frozen status via execution plan continuity
        """

    def validate_section(self, params: Dict = {"section_id": ""}):
        """
        Validate a research section by mutating the in-memory execution plan state.

        PURPOSE
        -------
        Marks a section as validate once it has been authored.
        A validate section is considered epistemically ready for freezing
        and must not be rewritten or re-analyzed within the same run.
        """

    def identify_required_sections(self):
        """
        Identify the authoritative section structure for a Deep Research run.

        PURPOSE
        -------
        Derives the complete, ordered list of research sections required to
        fully satisfy the user's request.

        This function establishes the **structural contract** of the research:
        • What analytical units must exist
        • In what order they must be addressed
        • What each section is responsible for conclusively answering

        Once identified, the section structure becomes authoritative and
        governs all subsequent execution, validation, and freezing logic.

        This function performs **structure identification only**.
        It does NOT plan execution.

        EXECUTION SEMANTICS
        -------------------
        • Runs exactly once per Deep Research run
        • Must complete before any section-level execution begins
        • Produces NO actions, NO lifecycle transitions, and NO execution steps
        • Output is immutable for the duration of the run

        INPUTS CONSIDERED
        -----------------
        • User query
        • Full conversation history
        • Enterprise and system context
        • Existing execution artifacts (if any)

        The function infers required sections even when:
        • The user uploads files
        • Requirements are spread across multiple messages
        • Sections are implied rather than explicitly listed

        WHAT THIS FUNCTION MUST NOT DO
        -------------------------------
        • Must NOT reference tools, actions, or functions
        • Must NOT describe data fetching or computation
        • Must NOT reason about execution order or lifecycle
        • Must NOT include validation, freezing, or confidence logic
        • Must NOT ask clarifying questions unless absolutely unavoidable

        OUTPUT
        ------
        Returns a structured JSON object conforming to
        `DEEP_RESEARCH_SECTION_PLAN_SCHEMA`, containing:

        • An ordered list of section definitions
        • Stable, descriptive section IDs
        • Clear objectives for each section
        • Declared analytical intent (deterministic / interpretive / mixed)
        • Explicit expectations for what constitutes section completion

        SIDE EFFECTS
        -------------
        • None
        • Does NOT mutate execution state
        • Does NOT persist data
        • Acts purely as a structural specification generator

        DOWNSTREAM CONTRACT
        -------------------
        • Section-level execution MUST NOT begin unless this function
        has successfully completed
        • The execution controller may only advance through sections
        defined here
        • No new sections may be introduced after this step
        """



